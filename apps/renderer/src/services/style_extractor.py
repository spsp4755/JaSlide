"""Extract reusable visual tokens from a PPTX without retaining its content."""

from collections import Counter
from io import BytesIO

from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn


def _solid_color(fill):
    if fill.type != MSO_FILL_TYPE.SOLID:
        return None
    try:
        value = str(fill.fore_color.rgb)
    except (AttributeError, TypeError, ValueError):
        return None
    return f"#{value}" if len(value) == 6 else None


def _frequent(values):
    counts = Counter(values)
    return sorted(counts, key=lambda value: (-counts[value], value))


def _font_name(run):
    properties = run._r.rPr
    east_asian = properties.find(qn("a:ea")) if properties is not None else None
    return east_asian.get("typeface") if east_asian is not None else run.font.name


def _html_layout(slide):
    text_shapes = [
        shape for shape in slide.shapes
        if shape.has_text_frame and shape.text_frame.text.strip()
    ]
    if not text_shapes:
        return None

    def font_size(shape):
        return max(
            (run.font.size.pt for paragraph in shape.text_frame.paragraphs for run in paragraph.runs if run.font.size),
            default=0,
        )

    slots = ("title", "body")
    shapes = sorted(text_shapes, key=lambda shape: (-font_size(shape), shape.top, shape.left))[: len(slots)]
    alignments = {PP_ALIGN.LEFT: "left", PP_ALIGN.CENTER: "center", PP_ALIGN.RIGHT: "right"}
    layout = []
    for index, shape in enumerate(shapes):
        attributes = [
            f'data-jaslide-slot="{slots[index]}"',
            f'data-x="{shape.left.inches:g}"',
            f'data-y="{shape.top.inches:g}"',
            f'data-w="{shape.width.inches:g}"',
            f'data-h="{shape.height.inches:g}"',
        ]
        size = font_size(shape)
        if size:
            attributes.append(f'data-font-size="{round(size):g}"')
        alignment = alignments.get(shape.text_frame.paragraphs[0].alignment)
        if alignment:
            attributes.append(f'data-align="{alignment}"')
        layout.append(f"<div {' '.join(attributes)}></div>")
    return "".join(layout)


def extract_template_tokens(pptx_bytes: bytes) -> dict:
    """Return deterministic color and font tokens; never expose slide text."""
    presentation = Presentation(BytesIO(pptx_bytes))
    backgrounds, fills, fonts = [], [], []

    for slide in presentation.slides:
        background = _solid_color(slide.background.fill)
        if background:
            backgrounds.append(background)
        for shape in slide.shapes:
            color = _solid_color(shape.fill)
            if color:
                fills.append(color)
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        font = _font_name(run)
                        if font:
                            fonts.append(font)

    colors = {}
    background = _frequent(backgrounds)
    if background:
        colors["background"] = background[0]
    palette = [color for color in _frequent(fills) if color != colors.get("background")]
    for key, color in zip(("primary", "secondary", "accent"), palette):
        colors[key] = color

    typography = {}
    font = _frequent(fonts)
    if font:
        typography = {"titleFont": font[0], "bodyFont": font[0]}
    config = {"colors": colors, "typography": typography}
    if presentation.slides:
        layout = _html_layout(presentation.slides[0])
        if layout:
            config["htmlTemplate"] = layout
    return config
