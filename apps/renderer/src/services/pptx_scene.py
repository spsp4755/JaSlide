"""Convert a PPTX deck into the canonical SlideScene model.

`pptx_to_html.py` produces the *rendering* — HTML the browser paints and, in
its `source.slides[].objects[]`, a flatter object map keyed for the live
canvas. This module produces the *storage source*: every shape as a scene
object with its full text runs (not just paragraph text), its exact OOXML
preset geometry name (not a lossy enum), and a `sourceRef` back to the shape
id everything else already keys edits by.

Reuses the low-level colour/font/geometry helpers from `pptx_to_html` rather
than re-deriving theme-colour and EMU-to-canvas-px math a second time —
importing them is worth the coupling to a private module; getting that math
subtly wrong twice is not.
"""

import copy
from io import BytesIO

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Picture

from .pptx_to_html import (
    CANVAS_HEIGHT,
    CANVAS_WIDTH,
    MAX_OBJECTS_PER_SLIDE,
    _color,
    _font_color,
    _font_name,
    _line_color,
    _line_width,
    _px,
    _pt_to_canvas_px,
)
from .theme_colors import theme_palette

# A shape with no preset geometry (a connector, a freeform) falls back to a
# plain rectangle — the same fidelity the HTML renderer already gives every
# shape it cannot draw a real silhouette for.
DEFAULT_SHAPE_KIND = "rect"


def _preset_geometry(shape) -> str | None:
    """The deck's own OOXML preset name (`prstGeom/@prst`), not an enum.

    `_preset_shape` on export feeds this string straight to
    `MSO_SHAPE.from_xml`, so keeping the deck's own spelling is what makes a
    re-exported shape come back as the same silhouette it started as.
    """
    element = shape._element.find(f".//{qn('a:prstGeom')}")
    return element.get("prst") if element is not None else None


def _run_dict(run, pt_to_px: float, palette: dict[str, str]) -> dict:
    return {
        "text": run.text,
        "bold": bool(run.font.bold),
        "italic": bool(run.font.italic),
        "underline": bool(run.font.underline),
        "color": _font_color(run.font.color, palette, run) or "#1A1A1A",
        "fontSize": round(run.font.size.pt) if run.font.size else None,
        "fontFamily": _font_name(run) or None,
    }


def _paragraph_dict(paragraph, pt_to_px: float, palette: dict[str, str]) -> dict:
    runs = [_run_dict(run, pt_to_px, palette) for run in paragraph.runs]
    # An empty line still has to round-trip as a paragraph, or a blank row in a
    # multi-line box silently disappears.
    if not runs:
        runs = [{"text": ""}]
    return {
        "runs": runs,
        "level": paragraph.level or 0,
        "align": {PP_ALIGN.CENTER: "center", PP_ALIGN.RIGHT: "right", PP_ALIGN.JUSTIFY: "justify"}.get(
            paragraph.alignment, "left"),
    }


def _text_frame_paragraphs(text_frame, pt_to_px: float, palette: dict[str, str]) -> list[dict]:
    return [_paragraph_dict(paragraph, pt_to_px, palette) for paragraph in text_frame.paragraphs]


def _table_cell_dict(cell, pt_to_px: float, palette: dict[str, str]) -> dict:
    result = {"paragraphs": _text_frame_paragraphs(cell.text_frame, pt_to_px, palette)}
    fill = _color(cell.fill, palette)
    if fill:
        result["fill"] = fill
    result["verticalAlign"] = {
        MSO_ANCHOR.MIDDLE: "middle", MSO_ANCHOR.BOTTOM: "bottom",
    }.get(cell.text_frame.vertical_anchor, "top")
    result["padding"] = {
        "top": round(cell.margin_top / 12700 * pt_to_px),
        "right": round(cell.margin_right / 12700 * pt_to_px),
        "bottom": round(cell.margin_bottom / 12700 * pt_to_px),
        "left": round(cell.margin_left / 12700 * pt_to_px),
    }
    if cell.is_spanned:
        result["spanned"] = True
        return result
    if cell.is_merge_origin:
        result["rowSpan"] = cell.span_height
        result["colSpan"] = cell.span_width
    borders = {}
    tc_pr = cell._tc.tcPr
    for xml_side, name in (("T", "top"), ("R", "right"), ("B", "bottom"), ("L", "left")):
        line = tc_pr.find(qn(f"a:ln{xml_side}")) if tc_pr is not None else None
        rgb = line.find(f".//{qn('a:srgbClr')}") if line is not None else None
        scheme = line.find(f".//{qn('a:schemeClr')}") if line is not None else None
        color = f"#{rgb.get('val').upper()}" if rgb is not None and rgb.get("val") else palette.get(scheme.get("val")) if scheme is not None else None
        if color:
            width = int(line.get("w", "12700"))
            borders[name] = {"color": color, "width": round(width / 12700)}
    if borders:
        result["border"] = borders
    return result


def pptx_to_scene(content: bytes) -> dict:
    """`{"slides": [SlideScene, ...]}` for every slide in the deck."""
    presentation = Presentation(BytesIO(content))
    pt_to_px = _pt_to_canvas_px(presentation)
    palette = theme_palette(presentation)
    slides = []

    def extract(shape, objects: list) -> None:
        if len(objects) >= MAX_OBJECTS_PER_SLIDE:
            return
        base = {
            "id": str(shape.shape_id),
            "x": _px(shape.left, presentation.slide_width, CANVAS_WIDTH),
            "y": _px(shape.top, presentation.slide_height, CANVAS_HEIGHT),
            "width": _px(shape.width, presentation.slide_width, CANVAS_WIDTH),
            "height": _px(shape.height, presentation.slide_height, CANVAS_HEIGHT),
            "rotation": shape.rotation or 0,
            "sourceRef": {"kind": "pptx-shape", "shapeId": str(shape.shape_id)},
        }

        if isinstance(shape, GroupShape):
            # A group's members carry their own absolute position already —
            # descend rather than treat the whole group as one opaque box.
            for member in shape.shapes:
                extract(member, objects)
        elif isinstance(shape, Picture):
            objects.append({**base, "type": "image", "src": ""})
        elif getattr(shape, "has_table", False):
            row_heights = [_px(row.height, presentation.slide_height, CANVAS_HEIGHT) for row in shape.table.rows]
            column_widths = [_px(column.width, presentation.slide_width, CANVAS_WIDTH) for column in shape.table.columns]
            objects.append({
                **base, "type": "table",
                # A table draws at the sum of its own column widths and row
                # heights, not the graphicFrame's stored extent — python-pptx
                # never recomputes that extent when a column is resized.
                "width": sum(column_widths) or base["width"],
                "height": sum(row_heights) or base["height"],
                "rowHeights": row_heights, "columnWidths": column_widths,
                "cells": [[_table_cell_dict(cell, pt_to_px, palette) for cell in row.cells] for row in shape.table.rows],
            })
        elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
            objects.append({
                **base, "type": "line",
                # ponytail: arrowheads are not read from the connector's XML yet
                # (headEnd/tailEnd) — every imported line comes back undecorated
                # until that's worth adding. Not a fidelity regression: the old
                # object map had no line kind at all, so this is a net gain.
                "lineStyle": "straightLine",
                "stroke": _line_color(shape, palette) or "#202124",
                "strokeWidth": _line_width(shape),
            })
        elif getattr(shape, "has_text_frame", False) and shape.text.strip():
            objects.append({
                **base, "type": "text",
                "paragraphs": _text_frame_paragraphs(shape.text_frame, pt_to_px, palette),
            })
        else:
            objects.append({
                **base, "type": "shape",
                "shape": _preset_geometry(shape) or DEFAULT_SHAPE_KIND,
                "fill": _color(getattr(shape, "fill", None), palette) or "#FFFFFF",
                "stroke": _line_color(shape, palette) or "#202124",
                "strokeWidth": _line_width(shape),
            })

    for slide in presentation.slides:
        objects: list[dict] = []
        for shape in slide.shapes:
            if len(objects) >= MAX_OBJECTS_PER_SLIDE:
                break
            try:
                extract(shape, objects)
            except Exception:  # noqa: BLE001
                # One exotic shape must not cost the user the whole slide.
                continue
        slides.append({"width": CANVAS_WIDTH, "height": CANVAS_HEIGHT, "objects": objects})

    return {"slides": slides}


def _table_cell_from_edit(cell) -> dict:
    """`_apply_native_edit` (pptx_generator.py) only ever reads a table edit's
    cells as plain strings, so that's the shape generation.service.ts writes —
    never the scene's own `{paragraphs: [...]}` cell. Normalize here, or
    SceneCanvas's table renderer crashes reading `.paragraphs` off a string."""
    if isinstance(cell, dict):
        return cell
    return {"paragraphs": [{"runs": [{"text": cell if isinstance(cell, str) else ""}], "level": 0, "align": "left"}]}


def _apply_edit_fields(object_: dict, edit: dict) -> None:
    """An existing scene object's fields, patched from one legacy `objectEdits`
    entry — everything `_apply_native_edit` (pptx_generator.py) reads."""
    for key, target in (("left", "x"), ("top", "y"), ("width", "width"), ("height", "height")):
        if isinstance(edit.get(key), (int, float)):
            object_[target] = edit[key]
    if isinstance(edit.get("rotation"), (int, float)):
        object_["rotation"] = edit["rotation"]
    if object_["type"] == "text" and isinstance(edit.get("paragraphs"), list):
        object_["paragraphs"] = edit["paragraphs"]
    elif object_["type"] == "text" and isinstance(edit.get("text"), str):
        # generation.service.ts also writes plain-text objectEdits the same way
        # _apply_native_edit (pptx_generator.py) reads them for a text shape —
        # {"text": "..."}, never the scene's own {paragraphs: [...]} shape. Mirror
        # that function's fallback here (one paragraph per "\n"-separated line,
        # each keeping the level/align of the pre-edit paragraph at the same
        # position, clamped to the last one) or the scene view keeps showing the
        # template's stale placeholder text instead of the generated one.
        existing = object_.get("paragraphs") or [{}]
        object_["paragraphs"] = [
            {
                "level": existing[min(index, len(existing) - 1)].get("level", 0),
                "align": existing[min(index, len(existing) - 1)].get("align", "left"),
                "runs": [{"text": line}],
            }
            for index, line in enumerate(edit["text"].split("\n"))
        ]
    elif object_["type"] == "table" and isinstance(edit.get("cells"), list):
        object_["cells"] = [
            [_table_cell_from_edit(cell) for cell in row] if isinstance(row, list) else row
            for row in edit["cells"]
        ]
    elif object_["type"] in ("shape", "line"):
        if isinstance(edit.get("fillColor"), str):
            object_["fill"] = edit["fillColor"]
        if isinstance(edit.get("lineColor"), str):
            object_["stroke"] = edit["lineColor"]
        if isinstance(edit.get("lineWidth"), (int, float)):
            object_["strokeWidth"] = edit["lineWidth"]


def _new_object_from_edit(object_id: str, edit: dict) -> dict | None:
    """A brand-new scene object for an `objectEdits` entry with no shape of its
    own in the source file yet — the mirror of `_inserted_object_edit` in
    `scene_to_pptx.py`. No `sourceRef`, same as every editor-inserted object."""
    base = {
        "id": object_id,
        "x": edit.get("left", 180), "y": edit.get("top", 180),
        "width": edit.get("width", 640), "height": edit.get("height", 100),
        "rotation": edit.get("rotation", 0),
    }
    if isinstance(edit.get("addText"), str):
        paragraphs = edit.get("paragraphs") or [{
            "runs": [{"text": edit.get("text", edit["addText"])}], "level": 0, "align": "left",
        }]
        return {**base, "type": "text", "paragraphs": paragraphs}
    if isinstance(edit.get("addTable"), dict):
        rows = max(1, int(edit["addTable"].get("rows", 3) or 3))
        columns = max(1, int(edit["addTable"].get("columns", 3) or 3))
        empty_cell = {"paragraphs": [{"runs": [{"text": ""}], "level": 0, "align": "left"}]}
        cells = edit.get("cells") or [[dict(empty_cell) for _ in range(columns)] for _ in range(rows)]
        return {
            **base, "type": "table",
            "rowHeights": [base["height"] / rows] * rows,
            "columnWidths": [base["width"] / columns] * columns,
            "cells": cells,
        }
    if isinstance(edit.get("addShape"), str):
        return {
            **base, "type": "shape", "shape": edit["addShape"],
            "fill": edit.get("fillColor", "#FFFFFF"), "stroke": edit.get("lineColor", "#202124"),
            "strokeWidth": edit.get("lineWidth", 2),
        }
    if isinstance(edit.get("addLine"), str):
        return {
            **base, "type": "line", "lineStyle": edit["addLine"],
            "stroke": edit.get("lineColor", "#202124"), "strokeWidth": edit.get("lineWidth", 2),
        }
    if isinstance(edit.get("imageData"), str):
        return {**base, "type": "image", "src": edit["imageData"]}
    return None


def apply_edits_to_scene(base_scene: dict, edits: list[dict]) -> dict:
    """Replay stored legacy `objectEdits` onto a freshly-imported `SlideScene`,
    producing the "current, already-edited" scene the editor should open
    with. The mirror of `scene_to_edits` in `scene_to_pptx.py` — both key
    objects by the same PPTX shape id, so this only needs pure data
    transforms, never a real PPTX render+reparse round trip."""
    objects = {object_["id"]: dict(object_) for object_ in base_scene["objects"]}
    order = [object_["id"] for object_ in base_scene["objects"]]
    for edit in edits:
        if not isinstance(edit, dict):
            continue
        object_id = edit.get("objectId")
        if not isinstance(object_id, str):
            continue
        if edit.get("delete") is True:
            objects.pop(object_id, None)
            if object_id in order:
                order.remove(object_id)
            continue
        duplicate_of = edit.get("duplicate")
        if isinstance(duplicate_of, str):
            source = objects.get(duplicate_of)
            if source is None:
                continue
            clone = copy.deepcopy(source)
            clone["id"] = object_id
            clone.pop("sourceRef", None)
            objects[object_id] = clone
            order.append(object_id)
            _apply_edit_fields(objects[object_id], edit)
            continue
        if object_id not in objects:
            inserted = _new_object_from_edit(object_id, edit)
            if inserted is not None:
                objects[object_id] = inserted
                order.append(object_id)
            continue
        _apply_edit_fields(objects[object_id], edit)
    return {**base_scene, "objects": [objects[object_id] for object_id in order if object_id in objects]}
