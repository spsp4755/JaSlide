"""Loss-minimizing PPTX-to-HTML conversion for template rendering."""

import base64
from html import escape
from io import BytesIO

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Picture

from .style_extractor import extract_template_tokens


CANVAS_WIDTH, CANVAS_HEIGHT = 1920, 1080
EMU_PER_INCH = 914400
# Used only when a shape states no size anywhere, matching python-pptx's own default.
DEFAULT_FONT_PT = 18


def _pt_to_canvas_px(presentation) -> float:
    """CSS px per point on the 1920x1080 canvas.

    Geometry is scaled from EMU into that canvas, but a font size arrives in
    points, and the two are not interchangeable: the canvas covers a 7.5in-tall
    slide, so it holds 144 px per inch and one point is two px. Emitting the point
    value as px — which this used to do — drew every deck's text at half size.
    """
    inches = (presentation.slide_height or 0) / EMU_PER_INCH
    return CANVAS_HEIGHT / inches / 72 if inches > 0 else 2.0
# A photo-heavy deck inlined every image as base64 and produced a 40MB config, which
# then has to live in a database row and load in the editor. For a PPTX template the
# blobs buy nothing — the preview is rendered from the original file and the editor
# works off the object map — so only small images (logos, icons) are worth embedding.
MAX_INLINE_IMAGE_BYTES = 256 * 1024
# One real deck had 5,888 shapes on a single slide once groups were unpacked. Drawing
# that many overlay handles locks the browser up, so keep the first slice in z-order.
MAX_OBJECTS_PER_SLIDE = 400
# Even at 256KB each, a 123-slide deck of photos still reached 26MB. Budget the whole
# deck: earlier slides keep their images, later ones fall back to a placeholder box.
MAX_INLINE_IMAGE_TOTAL_BYTES = 6 * 1024 * 1024


def _color(fill) -> str | None:
    try:
        value = fill.fore_color.rgb
        return f"#{value}" if value else None
    except (AttributeError, TypeError):
        return None


def _font_color(color) -> str | None:
    try:
        value = color.rgb
        return f"#{value}" if value else None
    except (AttributeError, TypeError):
        return None


def _font_name(run) -> str | None:
    properties = run._r.rPr
    east_asian = properties.find(qn("a:ea")) if properties is not None else None
    return east_asian.get("typeface") if east_asian is not None else run.font.name


def _line_width(shape) -> int:
    """Line width in px. A GroupShape has no `line` at all, and reading it unguarded
    raised AttributeError — which the API turned into "Invalid PPTX file", rejecting
    most real decks outright."""
    line = getattr(shape, "line", None)
    try:
        return max(1, round((line.width or 12700) / 12700)) if line is not None else 1
    except (AttributeError, TypeError):
        return 1


def _line_style(shape) -> str:
    try:
        if not shape.line.width:
            return ""
        color = _color(shape.line) or "#202124"
        width = max(1, round(shape.line.width / 12700))
        return f"border:{width}px solid {color}"
    except (AttributeError, TypeError):
        return ""


def _text_html(source, pt_to_px: float = 2.0) -> tuple[str, int, str | None]:
    """Markup, container font size in canvas px, and color for one text shape."""
    paragraphs = []
    # The container size is read back when this HTML is turned into a deck again, and
    # seeding it at 18 made it a floor: a 13pt caption came back as 18pt. Only fall
    # back to 18 when no run states a size at all.
    size = None
    color = None
    for paragraph in source.text_frame.paragraphs:
        runs = []
        for run in paragraph.runs:
            run_size = round(run.font.size.pt) if run.font.size else (size or DEFAULT_FONT_PT)
            run_color = _font_color(run.font.color) or color or "#1A1A1A"
            font_name = _font_name(run)
            size = run_size if size is None else max(size, run_size)
            color = color or run_color
            weight = "font-weight:700;" if run.font.bold else ""
            italic = "font-style:italic;" if run.font.italic else ""
            underline = "text-decoration:underline;" if run.font.underline else ""
            family = f'font-family:{escape(font_name, quote=True)};' if font_name else ""
            runs.append(f'<span style="font-size:{round(run_size * pt_to_px)}px;color:{run_color};{family}{weight}{italic}{underline}">{escape(run.text)}</span>')
        paragraphs.append("".join(runs) or escape(paragraph.text))
    return "<br>".join(paragraphs), round((size or DEFAULT_FONT_PT) * pt_to_px), color


def _text_style(shape) -> dict:
    """The object map's own formatting. Sizes here stay in points: python-pptx and the
    editor's on-slide text both want points, unlike the HTML above."""
    run = next((run for paragraph in shape.text_frame.paragraphs for run in paragraph.runs), None)
    if not run:
        return {}
    return {
        "fontSize": round(run.font.size.pt) if run.font.size else DEFAULT_FONT_PT,
        "color": _font_color(run.font.color) or "#1A1A1A",
        "fontFamily": _font_name(run) or "",
        "bold": bool(run.font.bold),
        "italic": bool(run.font.italic),
    }


def _text_align(shape) -> str:
    paragraph = next(iter(shape.text_frame.paragraphs), None)
    alignment = paragraph.alignment if paragraph is not None else None
    return {PP_ALIGN.CENTER: "center", PP_ALIGN.RIGHT: "right"}.get(alignment, "left")


def _table_html(shape, pt_to_px: float = 2.0) -> str:
    widths = [column.width for column in shape.table.columns]
    total_width = sum(widths) or 1
    rows = []
    for row in shape.table.rows:
        cells = []
        for index, cell in enumerate(row.cells):
            text, size, color = _text_html(cell, pt_to_px)
            fill = _color(cell.fill)
            width = round(widths[index] / total_width * 100, 1)
            surface = f"background:{fill};" if fill else ""
            # EMU / 12700 is points, not canvas px — scale it like every other length.
            cells.append(
                f'<td style="width:{width}%;height:{row.height / 12700 * pt_to_px:.1f}px;box-sizing:border-box;'
                f'border:1px solid #D1D5DB;padding:8px;vertical-align:middle;{surface}'
                f'font-size:{size}px;color:{color or "#1A1A1A"}">{text}</td>'
            )
        rows.append(f"<tr>{''.join(cells)}</tr>")
    return '<table style="width:100%;height:100%;border-collapse:collapse">' + "".join(rows) + "</table>"


def _px(value, total, canvas) -> int:
    return round(value / total * canvas)


def pptx_to_html(content: bytes) -> dict:
    presentation = Presentation(BytesIO(content))
    tokens = extract_template_tokens(content)
    pt_to_px = _pt_to_canvas_px(presentation)
    html_slides = []
    source_slides = []
    inlined = 0

    def extract(shape, objects: list, source_objects: list) -> None:
        """Turn one shape into its HTML and object-map entry, recursing into groups."""
        nonlocal inlined
        # Checked here, not only per top-level shape: one group can hold thousands.
        if len(source_objects) >= MAX_OBJECTS_PER_SLIDE:
            return
        left = _px(shape.left, presentation.slide_width, CANVAS_WIDTH)
        top = _px(shape.top, presentation.slide_height, CANVAS_HEIGHT)
        width = _px(shape.width, presentation.slide_width, CANVAS_WIDTH)
        height = _px(shape.height, presentation.slide_height, CANVAS_HEIGHT)
        position = f"position:absolute;left:{left}px;top:{top}px;width:{width}px;height:{height}px"
        source_object = {"id": str(shape.shape_id), "left": left, "top": top, "width": width, "height": height}
        # A picture dropped into a layout's placeholder reports shape_type
        # PLACEHOLDER, not PICTURE, so match on the class instead.
        if isinstance(shape, Picture):
            source_objects.append({**source_object, "kind": "image"})
            image = shape.image
            if len(image.blob) <= MAX_INLINE_IMAGE_BYTES and inlined + len(image.blob) <= MAX_INLINE_IMAGE_TOTAL_BYTES:
                inlined += len(image.blob)
                encoded = base64.b64encode(image.blob).decode("ascii")
                objects.append(f'<img data-object="true" data-object-type="image" src="data:{image.content_type};base64,{encoded}" style="{position}">')
            else:
                objects.append(f'<div data-object="true" data-object-type="image" style="{position};background:#E5E7EB"></div>')
        elif isinstance(shape, GroupShape):
            # Real decks group things constantly. Treating a group as one opaque box
            # made every diagram inside it uneditable, so descend into its members —
            # their own left/top are already absolute on the slide.
            for member in shape.shapes:
                extract(member, objects, source_objects)
        elif getattr(shape, "has_table", False):
            row_heights = [_px(row.height, presentation.slide_height, CANVAS_HEIGHT) for row in shape.table.rows]
            column_widths = [_px(column.width, presentation.slide_width, CANVAS_WIDTH) for column in shape.table.columns]
            # PowerPoint draws a table at the sum of its column widths and row
            # heights, not at the graphicFrame's stored extent — python-pptx never
            # recomputes that extent when a column is resized. Using the frame put
            # the preview and the editor's cell grid out of step with the real deck.
            width = sum(column_widths) or width
            height = sum(row_heights) or height
            position = f"position:absolute;left:{left}px;top:{top}px;width:{width}px;height:{height}px"
            source_objects.append({
                **source_object, "kind": "table", "width": width, "height": height,
                "cells": [[cell.text for cell in row.cells] for row in shape.table.rows],
                "rowHeights": row_heights, "columnWidths": column_widths,
            })
            objects.append(f'<div data-object="true" data-object-type="table" style="{position};box-sizing:border-box;overflow:hidden">{_table_html(shape, pt_to_px)}</div>')
        elif getattr(shape, "has_text_frame", False) and shape.text.strip():
            text, font_size, color = _text_html(shape, pt_to_px)
            source_objects.append({**source_object, "kind": "text", "text": shape.text, "align": _text_align(shape), "paragraphs": [{"text": paragraph.text, "level": paragraph.level} for paragraph in shape.text_frame.paragraphs], **_text_style(shape)})
            fill = _color(getattr(shape, "fill", None))
            surface = f"background:{fill};" if fill else ""
            objects.append(f'<div data-object="true" data-object-type="textbox" style="{position};box-sizing:border-box;overflow:hidden;{surface}{_line_style(shape)};font-size:{font_size}px;color:{color or "#1A1A1A"}">{text}</div>')
        else:
            fill = _color(getattr(shape, "fill", None))
            source_objects.append({**source_object, "kind": "shape", "fillColor": fill or "#FFFFFF",
                                   "lineColor": _color(getattr(shape, "line", None)) or "#202124",
                                   "lineWidth": _line_width(shape)})
            surface = f"background:{fill};" if fill else "background:transparent;"
            objects.append(f'<div data-object="true" data-object-type="shape" style="{position};box-sizing:border-box;{surface}{_line_style(shape)}"></div>')

    for slide in presentation.slides:
        background = _color(slide.background.fill) or "#FFFFFF"
        objects = []
        source_objects = []
        for shape in slide.shapes:
            if len(source_objects) >= MAX_OBJECTS_PER_SLIDE:
                break
            try:
                extract(shape, objects, source_objects)
            except Exception:  # noqa: BLE001
                # One exotic shape must not cost the user the whole upload: skip it and
                # keep the rest of the deck editable.
                continue
        html_slides.append(f'<div class="slide-container" style="position:relative;width:{CANVAS_WIDTH}px;height:{CANVAS_HEIGHT}px;overflow:hidden;background:{background}">{"".join(objects)}</div>')
        source_slides.append({"objects": source_objects})
    archive = {"slides": [f"slide-{index:02d}" for index in range(1, len(html_slides) + 1)], "canvas": {"width": CANVAS_WIDTH, "height": CANVAS_HEIGHT}}
    return {**tokens, "htmlSlides": html_slides, "htmlTemplate": html_slides[0] if html_slides else "", "archive": archive, "source": {"kind": "pptx", "slides": source_slides}}
