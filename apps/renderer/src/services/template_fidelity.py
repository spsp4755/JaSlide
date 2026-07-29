"""Report where a PPTX template's editable reconstruction falls short.

Not a pixel comparison — no rendering happens here. What is checkable without
one: which shapes this pipeline could not identify a real outline for (so they
render as a plain rectangle instead of their actual silhouette), which
connectors lost their arrow direction (documented as a known gap when
pptx_scene.py was built), and which fonts the deck names that are not
installed on this server, so the browser and LibreOffice will silently
substitute something else for them. An admin acting on this report knows
exactly what to expect from a re-extract, not just that "something" changed.
"""

import subprocess
from functools import lru_cache
from io import BytesIO

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.group import GroupShape
from pptx.shapes.picture import Picture

from .pptx_scene import _font_name, _preset_geometry


@lru_cache(maxsize=1)
def _installed_font_families() -> frozenset:
    """Every font family fontconfig on this server can resolve, read once per
    process — `fc-list` is a subprocess call, not worth repeating per report."""
    try:
        output = subprocess.run(
            ["fc-list", "--format=%{family}\n"],
            capture_output=True, text=True, timeout=5, check=False,
        ).stdout
    except (OSError, subprocess.SubprocessError):
        return frozenset()
    names = set()
    for line in output.splitlines():
        for name in line.split(","):
            name = name.strip()
            if name:
                names.add(name)
    return frozenset(names)


def _run_fonts(paragraphs) -> set:
    fonts = set()
    for paragraph in paragraphs:
        for run in paragraph.runs:
            name = _font_name(run)
            if name:
                fonts.add(name)
    return fonts


def fidelity_report(source_pptx: bytes) -> dict:
    """Extraction limits plus a small, actionable template support summary."""
    presentation = Presentation(BytesIO(source_pptx))
    degraded: list[dict] = []
    fonts_used: set = set()
    table_count = 0
    merged_cell_count = 0
    editable_object_count = 0

    def walk(shape) -> None:
        nonlocal table_count, merged_cell_count, editable_object_count
        if isinstance(shape, GroupShape):
            for member in shape.shapes:
                walk(member)
            return
        if isinstance(shape, Picture):
            editable_object_count += 1
            return
        if getattr(shape, "has_table", False):
            editable_object_count += 1
            table_count += 1
            for row in shape.table.rows:
                for cell in row.cells:
                    fonts_used.update(_run_fonts(cell.text_frame.paragraphs))
                    if cell.is_merge_origin:
                        merged_cell_count += 1
            return
        if getattr(shape, "has_text_frame", False) and shape.text.strip():
            editable_object_count += 1
            fonts_used.update(_run_fonts(shape.text_frame.paragraphs))
            return
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            editable_object_count += 1
            degraded.append({
                "objectId": str(shape.shape_id), "type": "line",
                "reason": "화살표 방향은 추출되지 않아 항상 직선으로 표시됩니다",
            })
            return
        if _preset_geometry(shape) is None:
            degraded.append({
                "objectId": str(shape.shape_id), "type": "shape",
                "reason": "원본 도형의 윤곽을 인식하지 못해 사각형으로 표시됩니다",
            })
        else:
            editable_object_count += 1

    for slide in presentation.slides:
        for shape in slide.shapes:
            try:
                walk(shape)
            except Exception:  # noqa: BLE001
                # One exotic shape must not cost the whole report.
                continue

    installed = _installed_font_families()
    missing_fonts = sorted(fonts_used - installed)

    return {
        "degradedObjects": degraded,
        "missingFontFamilies": missing_fonts,
        "summary": {
            "slides": len(presentation.slides),
            "editableObjects": editable_object_count,
            "tables": table_count,
            "mergedCells": merged_cell_count,
        },
    }
