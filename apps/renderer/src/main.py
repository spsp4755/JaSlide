"""
JaSlide Renderer - PPTX/PDF generation service
"""

from fastapi import FastAPI, File, HTTPException, UploadFile
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import Optional, List, Any
import io
import zipfile
from pptx import Presentation as PptxPresentation

from .generators.pptx_generator import PPTXGenerator
from .generators.pdf_exporter import PDFExporter
from .services.style_extractor import extract_template_tokens
from .services.html_template_archive import extract_html_template_archive

app = FastAPI(
    title="JaSlide Renderer",
    description="PPTX/PDF rendering service",
    version="0.1.0",
)

MAX_PPTX_ENTRIES = 500
MAX_PPTX_UNCOMPRESSED_BYTES = 100 * 1024 * 1024
REQUIRED_PPTX_ENTRIES = {"[Content_Types].xml", "_rels/.rels", "ppt/presentation.xml"}


class SlideContent(BaseModel):
    heading: Optional[str] = None
    subheading: Optional[str] = None
    body: Optional[str] = None
    bullets: Optional[List[dict]] = None
    image: Optional[dict] = None
    chart: Optional[dict] = None


class Slide(BaseModel):
    id: str
    order: int
    type: str
    title: Optional[str] = None
    content: dict
    layout: str = "center"
    notes: Optional[str] = None


class TemplateConfig(BaseModel):
    colors: Optional[dict] = None
    typography: Optional[dict] = None
    layouts: Optional[dict] = None
    backgrounds: Optional[dict] = None
    htmlTemplate: Optional[str] = None


class Template(BaseModel):
    id: str
    name: str
    config: Optional[TemplateConfig] = None


class Presentation(BaseModel):
    id: str
    title: str
    slides: List[Slide]
    template: Optional[Template] = None


class RenderRequest(BaseModel):
    presentation: Presentation


class PreviewRequest(BaseModel):
    presentation: Presentation
    slideIndex: int = 0


@app.get("/health")
async def health_check():
    return {"status": "ok", "service": "jaslide-renderer"}


def _is_safe_pptx_package(content: bytes) -> bool:
    try:
        with zipfile.ZipFile(io.BytesIO(content)) as package:
            entries = package.infolist()
            return (
                len(entries) <= MAX_PPTX_ENTRIES
                and sum(entry.file_size for entry in entries) <= MAX_PPTX_UNCOMPRESSED_BYTES
                and REQUIRED_PPTX_ENTRIES.issubset(package.namelist())
            )
    except (OSError, zipfile.BadZipFile, zipfile.LargeZipFile):
        return False


def _extract_pptx_content(content: bytes) -> dict:
    """Return only editable text grouped by slide; never retain media or formatting."""
    presentation = PptxPresentation(io.BytesIO(content))
    slides = []
    for number, slide in enumerate(presentation.slides, start=1):
        text_parts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                value = shape.text.strip()
                if value:
                    text_parts.append(value)
            elif getattr(shape, "has_table", False):
                rows = [" | ".join(cell.text.strip() for cell in row.cells if cell.text.strip()) for row in shape.table.rows]
                text_parts.extend(row for row in rows if row)
        slide_text = "\n".join(text_parts).strip()
        if slide_text:
            slides.append({
                "number": number,
                "title": text_parts[0].split("\n", 1)[0],
                "content": slide_text,
            })
    return {"content": "\n\n".join(slide["content"] for slide in slides), "slides": slides}


@app.post("/api/extract/style")
async def extract_style(file: UploadFile = File(...)):
    if (
        not (file.filename or "").lower().endswith(".pptx")
        or file.content_type
        != "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        raise HTTPException(status_code=400, detail="PPTX file required")
    content = await file.read()
    if not _is_safe_pptx_package(content):
        raise HTTPException(status_code=400, detail="Invalid PPTX package")
    try:
        config = extract_template_tokens(content)
    except Exception as error:
        raise HTTPException(status_code=400, detail="Invalid PPTX file") from error
    return {"config": config}


@app.post("/api/extract/html-template")
async def extract_html_template(file: UploadFile = File(...)):
    if not (file.filename or "").lower().endswith(".zip"):
        raise HTTPException(status_code=400, detail="ZIP file required")
    try:
        config = extract_html_template_archive(await file.read())
    except ValueError as error:
        raise HTTPException(status_code=400, detail=str(error)) from error
    return {"config": config}


@app.post("/api/extract/content")
async def extract_content(file: UploadFile = File(...)):
    if (
        not (file.filename or "").lower().endswith(".pptx")
        or file.content_type
        != "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        raise HTTPException(status_code=400, detail="PPTX file required")
    content = await file.read()
    if not _is_safe_pptx_package(content):
        raise HTTPException(status_code=400, detail="Invalid PPTX package")
    try:
        extracted = _extract_pptx_content(content)
    except Exception as error:
        raise HTTPException(status_code=400, detail="Invalid PPTX file") from error
    if not extracted["slides"]:
        raise HTTPException(status_code=400, detail="PPTX has no extractable content")
    return extracted


@app.post("/api/render/pptx")
async def render_pptx(request: RenderRequest):
    """Generate PPTX file from presentation data"""
    try:
        generator = PPTXGenerator(template_config=request.presentation.template)
        pptx_buffer = generator.generate(request.presentation)

        return StreamingResponse(
            io.BytesIO(pptx_buffer),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": f'attachment; filename="{request.presentation.title}.pptx"'
            },
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/render/pdf")
async def render_pdf(request: RenderRequest):
    """Generate PDF file from presentation data"""
    try:
        # First generate PPTX, then convert to PDF
        generator = PPTXGenerator(template_config=request.presentation.template)
        pptx_buffer = generator.generate(request.presentation)

        exporter = PDFExporter()
        pdf_buffer = exporter.convert_pptx_to_pdf(pptx_buffer)

        return StreamingResponse(
            io.BytesIO(pdf_buffer),
            media_type="application/pdf",
            headers={
                "Content-Disposition": f'attachment; filename="{request.presentation.title}.pdf"'
            },
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/render/preview")
async def render_preview(request: PreviewRequest):
    """Generate preview image for a specific slide"""
    try:
        generator = PPTXGenerator(template_config=request.presentation.template)
        pptx_buffer = generator.generate(request.presentation, request.slideIndex)
        preview_buffer = PDFExporter().convert_pptx_to_preview(pptx_buffer)

        return StreamingResponse(
            io.BytesIO(preview_buffer),
            media_type="image/png",
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(app, host="0.0.0.0", port=8000)
