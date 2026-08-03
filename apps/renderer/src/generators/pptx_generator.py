"""
PPTX Generator - Creates PowerPoint presentations using python-pptx
"""

from pptx import Presentation as PPTXPresentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# One mapping for every place an edit's alignment string reaches PP_ALIGN,
# so a fourth value (e.g. "justify") only needs adding once.
_PARAGRAPH_ALIGNMENTS = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from io import BytesIO
import base64
import copy
import logging
import math
from typing import Optional, Any
from ..services.html_template import extract_html_template_style, parse_html_layout, parse_html_objects
from ..services.html_renderer import render_slide_png

logger = logging.getLogger(__name__)

# A Hangul or CJK glyph occupies about a full em; Latin about half of one. Enough
# to tell "this obviously no longer fits" from "this still fits" without dragging
# in font metrics for every typeface a deck might use.
_WIDE_RANGES = ((0x1100, 0x11FF), (0x2E80, 0xA4CF), (0xAC00, 0xD7A3), (0xF900, 0xFAFF), (0xFF00, 0xFF60))


def _em_width(text: str) -> float:
    """Width of `text` in ems, counting East Asian glyphs as full width."""
    return sum(1.0 if any(low <= ord(ch) <= high for low, high in _WIDE_RANGES) else 0.5 for ch in text)


def fit_font_scale(lines: list[str], width_pt: float, height_pt: float, font_pt: float) -> float:
    """How much to shrink `font_pt` so `lines` still fit inside the box.

    Generated text is routinely longer than the text a template was built around,
    and PowerPoint's own answer to that is normAutofit — shrink rather than spill
    over whatever sits below. Returns a scale in (0, 1].
    """
    if not lines or width_pt <= 0 or height_pt <= 0 or font_pt <= 0:
        return 1.0
    for scale in (1.0, 0.9, 0.8, 0.7, 0.6, 0.5, 0.4, 0.3, 0.25):
        size = font_pt * scale
        wrapped = sum(max(1, math.ceil(_em_width(line) * size / width_pt)) for line in lines)
        if wrapped * size * 1.2 <= height_pt:
            return scale
    return 0.25


class PPTXGenerator:
    """Generate PPTX files from presentation data"""

    # Slide dimensions (16:9 aspect ratio)
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    # Default margins
    MARGIN_TOP = Inches(0.5)
    MARGIN_LEFT = Inches(0.5)
    MARGIN_RIGHT = Inches(0.5)
    MARGIN_BOTTOM = Inches(0.5)
    DEFAULT_COLORS = {"background": "#FFFFFF", "text": "#1E293B"}
    DEFAULT_FONT = "Noto Sans KR"

    def __init__(self, template_config: Optional[Any] = None, editable: bool = False):
        self.template_config = template_config
        self.editable = editable
        config = self._as_dict(getattr(template_config, "config", template_config))
        self.html_tokens, extracted_layout = extract_html_template_style(config.get("htmlTemplate", ""))
        self.html_slides = [slide for slide in config.get("htmlSlides", []) if isinstance(slide, str)]
        if not self.html_slides and isinstance(config.get("htmlTemplate"), str) and 'data-object="true"' in config["htmlTemplate"]:
            self.html_slides = [config["htmlTemplate"]]
        self.tokens = self._resolve_tokens(template_config)
        self.html_layout = parse_html_layout(config.get("htmlTemplate", "")) or extracted_layout
        self._reset_presentation()

    def _reset_presentation(self) -> None:
        self.prs = PPTXPresentation()
        self.prs.slide_width = self.SLIDE_WIDTH
        self.prs.slide_height = self.SLIDE_HEIGHT

    @staticmethod
    def _as_dict(value: Any) -> dict:
        if isinstance(value, dict):
            return value
        if hasattr(value, "model_dump"):
            return value.model_dump(exclude_none=True)
        return vars(value) if value is not None else {}

    @staticmethod
    def _rgb(value: Any, fallback: str) -> RGBColor:
        if not isinstance(value, str) or not value.startswith("#"):
            value = fallback
        value = value[1:]
        if len(value) != 6 or any(char not in "0123456789abcdefABCDEF" for char in value):
            value = fallback[1:]
        return RGBColor.from_string(value.upper())

    def _resolve_tokens(self, template: Any) -> dict:
        config = self._as_dict(getattr(template, "config", template))
        colors = self._as_dict(config.get("colors"))
        typography = self._as_dict(config.get("typography"))
        return {
            "background": self._rgb(colors.get("background") or self.html_tokens.get("background"), self.DEFAULT_COLORS["background"]),
            "text": self._rgb(colors.get("text") or self.html_tokens.get("text"), self.DEFAULT_COLORS["text"]),
            "title_font": typography.get("titleFont") or self.html_tokens.get("titleFont") or self.DEFAULT_FONT,
            "body_font": typography.get("bodyFont") or self.html_tokens.get("bodyFont") or self.DEFAULT_FONT,
        }

    def _apply_background(self, slide: Any) -> None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self.tokens["background"]

    def _layout(self, slot: str, defaults: dict) -> dict:
        return {**defaults, **self.html_layout.get(slot, {})}

    @staticmethod
    def _apply_alignment(paragraph: Any, value: Optional[str]) -> None:
        if value in _PARAGRAPH_ALIGNMENTS:
            paragraph.alignment = _PARAGRAPH_ALIGNMENTS[value]

    @staticmethod
    def _add_layout_textbox(slide: Any, layout: dict) -> Any:
        return slide.shapes.add_textbox(
            Inches(layout["x"]), Inches(layout["y"]), Inches(layout["w"]), Inches(layout["h"])
        )

    @staticmethod
    def _shrink_text_to_fit(shape: Any, default_pt: float = 18.0) -> None:
        """Ask PowerPoint to shrink this shape's text, and pre-compute the scale.

        Generated text is often longer than what the template was drawn around, so
        a title would run past its own box and over the table beneath it. A bare
        <a:normAutofit/> is enough for LibreOffice (which renders our previews) but
        PowerPoint keeps showing full size until the box is next edited, so write
        the fontScale it would have calculated too."""
        frame = getattr(shape, "text_frame", None)
        if frame is None or not shape.text.strip():
            return
        frame.word_wrap = True
        sizes = [run.font.size.pt for paragraph in frame.paragraphs for run in paragraph.runs if run.font.size]
        font_pt = max(sizes) if sizes else default_pt
        # Inset margins are part of the box the text has to live in.
        width_pt = max(1.0, (shape.width - frame.margin_left - frame.margin_right) / 12700)
        height_pt = max(1.0, (shape.height - frame.margin_top - frame.margin_bottom) / 12700)
        lines = [paragraph.text for paragraph in frame.paragraphs if paragraph.text]
        scale = fit_font_scale(lines, width_pt, height_pt, font_pt)
        body = frame._txBody.bodyPr
        for existing in body.findall(qn("a:normAutofit")) + body.findall(qn("a:spAutoFit")) + body.findall(qn("a:noAutofit")):
            body.remove(existing)
        autofit = OxmlElement("a:normAutofit")
        if scale < 1.0:
            autofit.set("fontScale", str(int(round(scale * 100000))))
            autofit.set("lnSpcReduction", "10000")
        body.append(autofit)

    def _style_paragraph(
        self, paragraph: Any, size: int, font: str, bold: bool = False, italic: bool = False
    ) -> None:
        for run in paragraph.runs:
            run.font.name = font
            r_pr = run._r.get_or_add_rPr()
            east_asian = r_pr.find(qn("a:ea"))
            if east_asian is None:
                east_asian = OxmlElement("a:ea")
                r_pr.append(east_asian)
            east_asian.set("typeface", font)
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = self.tokens["text"]

    def generate(self, presentation: Any, slide_index: Optional[int] = None) -> bytes:
        """Generate PPTX from presentation data"""
        config = self._as_dict(getattr(self.template_config, "config", self.template_config))
        source = config.get("sourcePptx")
        if isinstance(source, str) and source:
            self.prs = PPTXPresentation(BytesIO(base64.b64decode(source)))
            original_slide_count = len(self.prs.slides)
            slides_to_process = list(presentation.slides)
            if slide_index is not None:
                if slide_index < 0 or slide_index >= len(slides_to_process):
                    raise ValueError("Slide index is out of range")
                slides_to_process = [slides_to_process[slide_index]]

            used_slide_ids = set()
            kept_slide_ids = []
            for slide_data in slides_to_process:
                content = self._as_dict(getattr(slide_data, "content", {}))
                edits = content.get("objectEdits", [])
                if not original_slide_count:
                    continue
                # A template slide with nothing editable (a full-bleed image, say)
                # carries no edits, so fall back to the layout the generator picked
                # instead of silently rendering template slide 0 for every one.
                fallback = content.get("templateIndex") if isinstance(content.get("templateIndex"), int) else 0
                target = next(
                    (edit["slide"] for edit in edits if isinstance(edit, dict) and isinstance(edit.get("slide"), int)),
                    fallback,
                )
                target = max(0, min(target, original_slide_count - 1))
                original_slide = self.prs.slides[target]
                if original_slide.slide_id in used_slide_ids:
                    slide = self._clone_slide(target)
                else:
                    slide = original_slide
                    used_slide_ids.add(slide.slide_id)
                for edit in edits:
                    self._apply_native_edit(edit, slide)
                kept_slide_ids.append(slide.slide_id)

            if kept_slide_ids:
                # The export must contain exactly the generated slides, in their
                # generated order — not the untouched template deck with the
                # generated slides appended after it.
                self._keep_only_slides(kept_slide_ids)

            buffer = BytesIO()
            self.prs.save(buffer)
            return buffer.getvalue()
        self._reset_presentation()
        slides = list(enumerate(presentation.slides))
        if slide_index is not None:
            if slide_index < 0 or slide_index >= len(slides):
                raise ValueError("Slide index is out of range")
            slides = [slides[slide_index]]
        total_slides = len(presentation.slides)
        for template_index, slide_data in slides:
            self._add_slide(slide_data, template_index, total_slides)

        # Save to buffer
        buffer = BytesIO()
        self.prs.save(buffer)
        buffer.seek(0)
        return buffer.read()

    def _clone_slide(self, index: int) -> Any:
        """Duplicate an original template slide so a later generated slide's
        edits land on their own physical slide instead of overwriting an
        already-used one. Shape ids are preserved (deep-copied) so objectId
        based edit targeting still resolves on the clone."""
        source = self.prs.slides[index]
        dest = self.prs.slides.add_slide(source.slide_layout)
        for shape in list(dest.shapes):
            shape._element.getparent().remove(shape._element)
        for shape in source.shapes:
            new_element = copy.deepcopy(shape._element)
            for blip in new_element.iter(qn("a:blip")):
                old_rid = blip.get(qn("r:embed"))
                if old_rid and old_rid in source.part.rels:
                    rel = source.part.rels[old_rid]
                    new_rid = dest.part.rels._add_relationship(rel.reltype, rel.target_part)
                    blip.set(qn("r:embed"), new_rid)
            dest.shapes._spTree.append(new_element)
        return dest

    def _keep_only_slides(self, keep_slide_ids: list) -> None:
        keep = set(keep_slide_ids)
        sld_id_lst = self.prs.slides._sldIdLst
        for sld in list(sld_id_lst):
            if int(sld.get("id")) not in keep:
                self.prs.part.drop_rel(sld.get(qn("r:id")))
                sld_id_lst.remove(sld)
        # Re-appending an existing child moves it, so this reorders the deck to
        # match the generated slide order.
        remaining = {int(sld.get("id")): sld for sld in sld_id_lst}
        for slide_id in keep_slide_ids:
            if slide_id in remaining:
                sld_id_lst.append(remaining[slide_id])

    @staticmethod
    def _preset_shape(kind: str) -> Any:
        """Resolve an OOXML preset geometry name to its MSO_SHAPE member.

        The editor sends preset names straight from the shape picker, so all 180-odd
        presets work. The old hand-written lookup covered 20 of them and silently
        exported everything else as a rectangle."""
        try:
            return MSO_SHAPE.from_xml(kind)
        except (KeyError, ValueError):
            return MSO_SHAPE.RECTANGLE

    @staticmethod
    def _style_connector(shape: Any, kind: str) -> None:
        """Add the arrowheads and dash pattern the picker's icon promised.

        python-pptx has no API for either, so write the <a:ln> children directly."""
        line = shape.line._get_or_add_ln()
        for name, present in (("a:headEnd", kind.startswith("double")), ("a:tailEnd", "rrow" in kind)):
            if present:
                end = OxmlElement(name)
                end.set("type", "triangle")
                line.append(end)
        dash = {"dashedLine": "dash", "dottedLine": "sysDot"}.get(kind)
        if dash:
            element = OxmlElement("a:prstDash")
            element.set("val", dash)
            line.append(element)

    @staticmethod
    def _safe_rgb(color: Any):
        """python-pptx raises AttributeError reading .rgb off a theme/scheme
        color (common in real-world decks); treat that as "no explicit color"
        instead of crashing the whole render."""
        try:
            return color.rgb if color.type is not None else None
        except AttributeError:
            return None

    _BULLET_MARKER_CHARS = "-–—•·∙‣▪▫◦*"

    @classmethod
    def _strip_leading_bullet_marker(cls, text: str) -> str:
        """Drop a literal bullet character the model wrote at the start of a
        line — the template's own a:buChar already draws one; writing both
        doubles it up."""
        stripped = text.lstrip(" \t")
        indent = text[: len(text) - len(stripped)]
        marker_end = 0
        while marker_end < len(stripped) and stripped[marker_end] in cls._BULLET_MARKER_CHARS:
            marker_end += 1
        if marker_end == 0:
            return text
        rest = stripped[marker_end:]
        if rest and not rest[0] in " \t":
            return text  # not actually a marker followed by a space — e.g. "1-2" or a real hyphenated word
        return indent + rest.lstrip(" \t")

    @staticmethod
    def _run_has_explicit_style(run_item: Optional[dict]) -> bool:
        if not run_item:
            return False
        if any(run_item.get(key) is True for key in ("bold", "italic", "underline")):
            return True
        return any(run_item.get(key) for key in ("color", "fontSize", "fontFamily"))

    @staticmethod
    def _set_first_run_text(paragraph_element, text: str) -> None:
        runs = paragraph_element.findall(qn("a:r"))
        for extra in runs[1:]:
            paragraph_element.remove(extra)
        if not runs:
            return
        text_element = runs[0].find(qn("a:t"))
        if text_element is None:
            return
        text_element.text = text

    def _write_paragraphs(self, frame: Any, paragraphs: list) -> None:
        # A real template's first paragraph in a list frequently carries
        # different formatting from the rest (e.g. a distinct "header"
        # bullet/indent) even though python-pptx reports the same numeric
        # level for all of them. Keeping only the FIRST prototype found per
        # level -- instead of collecting every distinctly-styled paragraph
        # and cycling through them by call order -- guarantees every
        # generated line at a given level renders identically, matching
        # what "level" is supposed to mean visually.
        prototype_by_level: dict[int, Any] = {}
        for paragraph in frame.paragraphs:
            if not paragraph.runs:
                continue  # a run-less prototype has nowhere to put the generated text
            level = paragraph.level or 0
            if level not in prototype_by_level:
                prototype_by_level[level] = copy.deepcopy(paragraph._p)

        # frame.clear() special-cases the first paragraph (keeps its <a:pPr>,
        # only strips runs) instead of removing it — remove every existing
        # paragraph ourselves so nothing is left over before we start writing.
        for existing in list(frame._txBody.p_lst):
            frame._txBody.remove(existing)

        def pick_prototype(level: int):
            if not prototype_by_level:
                return None
            if level in prototype_by_level:
                return prototype_by_level[level]
            nearest = min(prototype_by_level, key=lambda existing_level: abs(existing_level - level))
            return prototype_by_level[nearest]

        for item in paragraphs:
            if not isinstance(item, dict):
                continue
            runs = item.get("runs")
            level = max(0, item["level"]) if isinstance(item.get("level"), int) else 0
            single_run = runs[0] if isinstance(runs, list) and len(runs) == 1 and isinstance(runs[0], dict) else None
            is_simple = (not isinstance(runs, list) or len(runs) <= 1) and not self._run_has_explicit_style(single_run)

            if is_simple:
                prototype = pick_prototype(level)
                if prototype is not None:
                    # The edit's own "align" is not applied here on purpose: every
                    # AI-generated line currently carries a uniform "left" default
                    # regardless of the template's real alignment, so honoring it
                    # would silently override a deliberately centered/right-aligned
                    # template paragraph with a meaningless value. The template's
                    # own alignment wins on this path; only the from-scratch branch
                    # below (explicit per-run styling) applies a requested align.
                    text = str(single_run.get("text", "")) if single_run else str(item.get("text", ""))
                    text = self._strip_leading_bullet_marker(text)
                    clone = copy.deepcopy(prototype)
                    self._set_first_run_text(clone, text)
                    frame._txBody.append(clone)
                    continue

            # Explicit per-run styling (chat-edit character formatting) or a
            # template with no paragraphs at all to clone from — build fresh,
            # exactly as before.
            paragraph = frame.add_paragraph()
            if isinstance(item.get("level"), int):
                paragraph.level = level
            if isinstance(item.get("align"), str):
                alignment = _PARAGRAPH_ALIGNMENTS.get(item["align"])
                if alignment is not None:
                    paragraph.alignment = alignment
            if isinstance(runs, list) and runs:
                for run_item in runs:
                    if not isinstance(run_item, dict):
                        continue
                    run = paragraph.add_run()
                    run.text = str(run_item.get("text", ""))
                    if isinstance(run_item.get("bold"), bool):
                        run.font.bold = run_item["bold"]
                    if isinstance(run_item.get("italic"), bool):
                        run.font.italic = run_item["italic"]
                    if isinstance(run_item.get("underline"), bool):
                        run.font.underline = run_item["underline"]
                    if isinstance(run_item.get("color"), str) and len(run_item["color"].lstrip("#")) == 6:
                        run.font.color.rgb = RGBColor.from_string(run_item["color"].lstrip("#").upper())
                    if isinstance(run_item.get("fontSize"), (int, float)):
                        run.font.size = Pt(run_item["fontSize"])
                    if isinstance(run_item.get("fontFamily"), str):
                        run.font.name = run_item["fontFamily"]
            else:
                paragraph.text = str(item.get("text", ""))

        # A text frame must always have at least one <a:p> — OOXML's CT_TextBody
        # requires minOccurs="1". An edit with an empty paragraphs list must not
        # leave the shape with zero paragraphs, or PowerPoint treats the file as
        # needing repair. If nothing was appended, add a single empty paragraph.
        if not frame._txBody.p_lst:
            frame._txBody.add_p()

    def _apply_native_edit(self, edit: dict, slide: Any) -> None:
        if not isinstance(edit, dict):
            return
        shape = next((item for item in slide.shapes if str(item.shape_id) == str(edit.get("objectId"))), None)
        if edit.get("delete") is True:
            if shape:
                shape._element.getparent().remove(shape._element)
            return
        left = int((edit.get("left", 180)) * self.prs.slide_width / 1920)
        top = int((edit.get("top", 180)) * self.prs.slide_height / 1080)
        width = int((edit.get("width", 640)) * self.prs.slide_width / 1920)
        height = int((edit.get("height", 100)) * self.prs.slide_height / 1080)
        if isinstance(edit.get("duplicate"), str):
            # Ctrl+D on a template object. Copy the element wholesale so the clone keeps
            # its geometry, fill and text, then let the rest of this edit reposition it.
            original = next((item for item in slide.shapes if str(item.shape_id) == edit["duplicate"]), None)
            if original is not None:
                clone = copy.deepcopy(original._element)
                for properties in clone.iter(qn("p:cNvPr")):
                    properties.set("id", str(max((item.shape_id for item in slide.shapes), default=1) + 1))
                slide.shapes._spTree.append(clone)
                shape = next((item for item in slide.shapes if item._element is clone), None)
        elif isinstance(edit.get("addTable"), dict):
            # Inserting a table was impossible on a PPTX slide: the toolbar button only
            # rewrote `content.html`, which a PPTX-backed slide does not have.
            rows = max(1, min(int(edit["addTable"].get("rows", 3) or 3), 30))
            columns = max(1, min(int(edit["addTable"].get("columns", 3) or 3), 20))
            shape = slide.shapes.add_table(rows, columns, left, top, width, height)
        elif isinstance(edit.get("addShape"), str):
            shape = slide.shapes.add_shape(self._preset_shape(edit["addShape"]), left, top, width, height)
        elif isinstance(edit.get("addLine"), str):
            kind = edit["addLine"]
            connector = MSO_CONNECTOR.CURVE if kind.startswith("curved") else MSO_CONNECTOR.ELBOW if "elbow" in kind else MSO_CONNECTOR.STRAIGHT
            shape = slide.shapes.add_connector(connector, left, top + height // 2, left + width, top + height // 2)
            self._style_connector(shape, kind)
        elif isinstance(edit.get("addText"), str):
            shape = slide.shapes.add_textbox(left, top, width, height)
            shape.text = edit.get("text", edit["addText"])
        if isinstance(edit.get("imageData"), str) and "," in edit["imageData"]:
            try:
                image = base64.b64decode(edit["imageData"].split(",", 1)[1])
                left = int((edit.get("left", 180)) * self.prs.slide_width / 1920)
                top = int((edit.get("top", 180)) * self.prs.slide_height / 1080)
                width = int((edit.get("width", 640)) * self.prs.slide_width / 1920)
                height = int((edit.get("height", 360)) * self.prs.slide_height / 1080)
                slide.shapes.add_picture(BytesIO(image), left, top, width, height)
            except Exception:
                # Losing one pasted image should not fail the whole deck, but a
                # silent drop left no way to tell an unsupported file apart from
                # an image that never arrived.
                logger.warning("Skipped an image edit that could not be decoded", exc_info=True)
            return
        if not shape:
            return
        if getattr(shape, "has_text_frame", False) and isinstance(edit.get("paragraphs"), list):
            self._write_paragraphs(shape.text_frame, edit["paragraphs"])
        elif isinstance(edit.get("text"), str) and getattr(shape, "has_text_frame", False):
            levels = [paragraph.level for paragraph in shape.text_frame.paragraphs]
            alignments = [paragraph.alignment for paragraph in shape.text_frame.paragraphs]
            shape.text = edit["text"]
            for index, paragraph in enumerate(shape.text_frame.paragraphs):
                paragraph.level = levels[min(index, len(levels) - 1)] if levels else 0
                paragraph.alignment = alignments[min(index, len(alignments) - 1)] if alignments else None
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if isinstance(edit.get("fontFamily"), str):
                        run.font.name = edit["fontFamily"]
                    if isinstance(edit.get("fontSize"), (int, float)):
                        run.font.size = Pt(edit["fontSize"])
                    if isinstance(edit.get("color"), str) and len(edit["color"].lstrip("#")) == 6:
                        run.font.color.rgb = RGBColor.from_string(edit["color"].lstrip("#").upper())
                    if isinstance(edit.get("bold"), bool):
                        run.font.bold = edit["bold"]
                    if isinstance(edit.get("italic"), bool):
                        run.font.italic = edit["italic"]
                    if isinstance(edit.get("underline"), bool):
                        run.font.underline = edit["underline"]
            # Alignment is a paragraph property, not a run one — setting it per run
            # silently did nothing, so the toolbar's align buttons never took.
            if isinstance(edit.get("align"), str):
                alignment = _PARAGRAPH_ALIGNMENTS.get(edit["align"])
                if alignment is not None:
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph.alignment = alignment
        for key, size in (("left", self.prs.slide_width), ("width", self.prs.slide_width), ("top", self.prs.slide_height), ("height", self.prs.slide_height)):
            if isinstance(edit.get(key), (int, float)):
                setattr(shape, key, int(edit[key] * size / (1920 if key in ("left", "width") else 1080)))
        if isinstance(edit.get("fillColor"), str) and len(edit["fillColor"].lstrip("#")) == 6:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor.from_string(edit["fillColor"].lstrip("#").upper())
        if isinstance(edit.get("lineColor"), str) and len(edit["lineColor"].lstrip("#")) == 6:
            shape.line.color.rgb = RGBColor.from_string(edit["lineColor"].lstrip("#").upper())
        if isinstance(edit.get("lineWidth"), (int, float)):
            shape.line.width = int(max(0, edit["lineWidth"]) * 12700)
        # An inserted arrow that cannot be turned only ever points right.
        if isinstance(edit.get("rotation"), (int, float)):
            try:
                shape.rotation = float(edit["rotation"]) % 360
            except (AttributeError, ValueError):
                pass

        # Z-order is the shape's position in the slide's shape tree. Without this an
        # object could be moved but never pulled out from under an overlapping one.
        if edit.get("order") in ("front", "back"):
            tree = shape._element.getparent()
            tree.remove(shape._element)
            if edit["order"] == "front":
                tree.append(shape._element)
            else:
                # index 2 skips the tree's required non-shape children (nvGrpSpPr, grpSpPr).
                tree.insert(2, shape._element)
        cells = edit.get("cells")
        if isinstance(cells, list) and getattr(shape, "has_table", False):
            for row_index, row in enumerate(cells):
                if not isinstance(row, list) or row_index >= len(shape.table.rows):
                    continue
                for column_index, value in enumerate(row):
                    if column_index < len(shape.table.columns) and isinstance(value, str):
                        cell = shape.table.cell(row_index, column_index)
                        source = cell.text_frame.paragraphs
                        source_runs = [{"name": run.font.name, "size": run.font.size, "bold": run.font.bold, "italic": run.font.italic, "color": self._safe_rgb(run.font.color)} if (run := (paragraph.runs[0] if paragraph.runs else None)) else None for paragraph in source]
                        levels = [paragraph.level for paragraph in source]
                        alignments = [paragraph.alignment for paragraph in source]
                        cell.text = value
                        for index, paragraph in enumerate(cell.text_frame.paragraphs):
                            paragraph.level = levels[min(index, len(levels) - 1)] if levels else 0
                            paragraph.alignment = alignments[min(index, len(alignments) - 1)] if alignments else None
                            source_run = source_runs[min(index, len(source_runs) - 1)] if source_runs else None
                            if source_run and paragraph.runs:
                                run = paragraph.runs[0]
                                run.font.name = source_run["name"]
                                run.font.size = source_run["size"]
                                run.font.bold = source_run["bold"]
                                run.font.italic = source_run["italic"]
                                if source_run["color"]:
                                    run.font.color.rgb = source_run["color"]
                    elif column_index < len(shape.table.columns) and isinstance(value, dict) and isinstance(value.get("paragraphs"), list):
                        self._write_paragraphs(shape.table.cell(row_index, column_index).text_frame, value["paragraphs"])
        # Last, so it measures the text that actually ended up in the shape, at
        # whatever size and box this edit left behind.
        if getattr(shape, "has_text_frame", False):
            self._shrink_text_to_fit(shape)

    def _add_slide(self, slide_data: Any, template_index: int = 0, total_slides: int = 1):
        """Add a slide based on its type"""
        content = self._as_dict(getattr(slide_data, "content", {}))
        # The scene editor is the same authoritative source the API's GetScene reads
        # first (service.go checks content["scene"] before objectEdits/html too) — a
        # manual edit here must reach the export instead of being silently overwritten
        # by a fresh regeneration from content.heading/bullets.
        scene = content.get("scene")
        if isinstance(scene, dict) and scene.get("objects"):
            self._add_scene_slide(scene)
            return
        # A rendered screenshot matches the HTML exactly but exports one flat picture per
        # slide: nothing is selectable once the file leaves JaSlide. `editable` trades a
        # little fidelity for a deck the recipient can actually revise.
        if isinstance(content.get("html"), str) and content["html"].strip():
            if not self.editable:
                self._add_html_image_slide(content["html"])
                return
            # Transcribe this slide's own HTML, so the editable export carries the text
            # that is actually on the slide rather than the template's empty layout.
            objects = parse_html_objects(content["html"])
            if objects:
                self._add_editable_html_slide(objects)
                return
        if self.html_slides:
            selected_index = self._template_index(slide_data, template_index, total_slides)
            objects = parse_html_objects(self.html_slides[selected_index])
            if objects:
                self._add_html_template_slide(slide_data, selected_index, objects)
                return
            # ponytail: the uploaded deck isn't in JaSlide's data-object markup, so there's
            # nothing to place absolutely. Fall through to the generic layout below, which
            # already picks up this template's extracted background/font tokens (self.tokens)
            # instead of emitting a blank slide.
        slide_type = slide_data.type.upper()
        content = slide_data.content

        # ponytail: small models often ignore an outline's requested slide type
        # (e.g. picking CONTENT instead of the guided TWO_COLUMN) while still
        # correctly producing a valid two-column "columns" content shape. Trust
        # the content shape over the unreliable type label when it's present.
        columns = content.get("columns") if isinstance(content, dict) else None
        # A table/chart shape always wins over columns: parseSlideContent (Go)
        # attaches "columns" unconditionally regardless of slideType, and a
        # TABLE/CHART slide always carries a real or placeholder table/chart —
        # without this guard, a model that (sloppily) emits both shapes on a
        # TABLE/CHART slide would have its table/chart silently dropped.
        has_columns = (
            isinstance(columns, list) and len(columns) == 2
            and all(isinstance(item, dict) for item in columns)
            and not content.get("table") and not content.get("chart")
        )
        # table/chart/columns all outrank the four layouts below, for the same
        # reason table/chart outranks columns above: parseSlideContent (Go)
        # attaches each of these shapes unconditionally whenever it validates,
        # regardless of slideType, so a model that emits more than one shape on
        # the same slide must not have the "later" shape silently dropped.
        timeline = content.get("timeline") if isinstance(content, dict) else None
        has_timeline = (
            isinstance(timeline, dict) and isinstance(timeline.get("items"), list)
            and 3 <= len(timeline["items"]) <= 8
            and not content.get("table") and not content.get("chart") and not has_columns
        )
        process = content.get("process") if isinstance(content, dict) else None
        has_process = (
            isinstance(process, dict) and isinstance(process.get("steps"), list)
            and 2 <= len(process["steps"]) <= 6
            and not content.get("table") and not content.get("chart") and not has_columns
        )
        comparison = content.get("comparison") if isinstance(content, dict) else None
        has_comparison = (
            isinstance(comparison, dict) and isinstance(comparison.get("left"), dict) and isinstance(comparison.get("right"), dict)
            and not content.get("table") and not content.get("chart") and not has_columns
        )
        metrics = content.get("metrics") if isinstance(content, dict) else None
        has_kpi = (
            isinstance(metrics, dict) and isinstance(metrics.get("metrics"), list)
            and 2 <= len(metrics["metrics"]) <= 6
            and not content.get("table") and not content.get("chart") and not has_columns
        )

        if slide_type == "TITLE":
            self._add_title_slide(slide_data)
        elif has_columns or slide_type == "TWO_COLUMN":
            self._add_two_column_slide(slide_data)
        elif has_timeline:
            self._add_timeline_slide(slide_data)
        elif has_process:
            self._add_process_slide(slide_data)
        elif has_comparison:
            self._add_comparison_slide(slide_data)
        elif has_kpi:
            self._add_kpi_slide(slide_data)
        elif slide_type == "CONTENT":
            self._add_content_slide(slide_data)
        elif slide_type == "BULLET_LIST":
            self._add_bullet_slide(slide_data)
        elif slide_type == "QUOTE":
            self._add_quote_slide(slide_data)
        elif slide_type == "SECTION_HEADER":
            self._add_section_header_slide(slide_data)
        else:
            self._add_content_slide(slide_data)

    def _add_scene_slide(self, scene: dict) -> None:
        """A slide with no PPTX/HTML source, edited directly as a SlideScene. None of
        its objects have a source PPTX shape, so every one is placed with the same
        insert-only edit path scene_to_pptx uses for brand-new shapes."""
        from ..services.scene_to_pptx import scene_to_edits
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        for edit in scene_to_edits(scene):
            self._apply_native_edit(edit, slide)

    def _add_editable_html_slide(self, objects: list[dict]) -> None:
        """Place every HTML object as a real shape or textbox, keeping its own text.

        Back-to-front, so an object that sat on top in the HTML still sits on top.

        ponytail: an HTML <table> arrives as one flattened string, so a table-heavy
        slide loses its grid here — the default image export still renders it exactly.
        Carry cells through parse_html_objects and emit add_table if that matters."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        full_bleed = next((item for item in objects if item["x"] <= 0.05 and item["y"] <= 0.05
                           and item["w"] >= 13.3 and item["h"] >= 7.4 and item["background"]), None)
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = self._rgb(full_bleed["background"], self.DEFAULT_COLORS["background"]) if full_bleed else self.tokens["background"]

        for item in objects:
            if item is full_bleed:
                continue
            # A grid must come out a grid: flattening it into one string turned an ASR
            # table into a run-on line of headers with the rows gone.
            if item.get("cells"):
                self._add_editable_html_table(slide, item)
                continue
            if item["background"]:
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(item["x"]), Inches(item["y"]), Inches(item["w"]), Inches(item["h"]))
                shape.fill.solid()
                shape.fill.fore_color.rgb = self._rgb(item["background"], self.DEFAULT_COLORS["background"])
                shape.line.fill.background()
            if not item["text"]:
                continue
            box = self._add_layout_textbox(slide, {"x": item["x"] + 0.06, "y": item["y"] + 0.04,
                                                  "w": max(item["w"] - 0.12, 0.2), "h": max(item["h"] - 0.08, 0.2)})
            box.text_frame.word_wrap = True
            paragraph = box.text_frame.paragraphs[0]
            paragraph.text = item["text"]
            self._style_paragraph(paragraph, item["fontSize"], item["font"] or self.tokens["body_font"], bold=item["bold"])
            # A caption usually has no background of its own — what decides its colour is
            # the panel behind it. Reading only its own background put black on black.
            behind = next((other for other in reversed(objects[:objects.index(item)])
                           if other["background"]
                           and other["x"] <= item["x"] and other["y"] <= item["y"]
                           and other["x"] + other["w"] >= item["x"] + item["w"]
                           and other["y"] + other["h"] >= item["y"] + item["h"]), full_bleed)
            backdrop = item["background"] or (behind["background"] if behind else "")
            for run in paragraph.runs:
                run.font.color.rgb = self._rgb(item["color"] or ("#FFFFFF" if self._is_dark(backdrop) else "#1A1A1A"), self.DEFAULT_COLORS["text"])
            self._apply_alignment(paragraph, item["align"])

    def _add_editable_html_table(self, slide: Any, item: dict) -> None:
        """Rebuild an HTML table as a real PPTX table, so its rows survive the export."""
        rows = [row for row in item["cells"] if any(cell for cell in row)]
        columns = max((len(row) for row in rows), default=0)
        if not rows or not columns:
            return
        graphic = slide.shapes.add_table(len(rows), columns, Inches(item["x"]), Inches(item["y"]),
                                        Inches(item["w"]), Inches(item["h"]))
        table = graphic.table
        for row_index, row in enumerate(rows):
            for column_index in range(columns):
                cell = table.cell(row_index, column_index)
                cell.text = row[column_index] if column_index < len(row) else ""
                paragraph = cell.text_frame.paragraphs[0]
                # The first row of an HTML table is its header often enough to be worth
                # carrying over; a reader can undo one bold row far more easily than
                # rebuilding a grid.
                self._style_paragraph(paragraph, max(9, min(item["fontSize"], 16)),
                                      item["font"] or self.tokens["body_font"], bold=row_index == 0)

    def _add_html_image_slide(self, html: str) -> None:
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.shapes.add_picture(BytesIO(render_slide_png(html)), 0, 0, self.SLIDE_WIDTH, self.SLIDE_HEIGHT)

    def _add_html_template_slide(self, slide_data: Any, selected_index: int, objects: list[dict]):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        background = next((item["background"] for item in objects if item["x"] == 0 and item["y"] == 0 and item["w"] >= 13.3 and item["h"] >= 7.4 and item["background"]), None)
        if background:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = self._rgb(background, self.DEFAULT_COLORS["background"])
        else:
            self._apply_background(slide)

        content = slide_data.content or {}
        content_slots = []
        for item in objects:
            if item["type"] == "shape" and item["background"] and not (item["x"] == 0 and item["y"] == 0 and item["w"] >= 13.3 and item["h"] >= 7.4):
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(item["x"]), Inches(item["y"]), Inches(item["w"]), Inches(item["h"]))
                shape.fill.solid()
                shape.fill.fore_color.rgb = self._rgb(item["background"], self.DEFAULT_COLORS["background"])
                shape.line.fill.background()
                if item["w"] >= 2 and item["h"] >= 0.7:
                    content_slots.append(item)

        if self._add_semantic_html_layout(slide, self._template_name(selected_index), slide_data, content):
            return

        textboxes = sorted((item for item in objects if item["type"] == "textbox"), key=lambda item: item["fontSize"], reverse=True)
        generated_text = [content.get("heading", slide_data.title or "")]
        for item in textboxes:
            box = self._add_layout_textbox(slide, item)
            text = generated_text[textboxes.index(item)] if item in textboxes[:1] else ""
            box.text_frame.paragraphs[0].text = text
            self._style_paragraph(box.text_frame.paragraphs[0], item["fontSize"], item["font"] or self.tokens["body_font"], bold=item["bold"])
            box.text_frame.paragraphs[0].font.color.rgb = self._rgb(item["color"], self.DEFAULT_COLORS["text"])
            self._apply_alignment(box.text_frame.paragraphs[0], item["align"])

        if str(getattr(slide_data, "type", "")).upper() == "CHART" and self._add_chart(slide, content, content_slots):
            return

        if str(getattr(slide_data, "type", "")).upper() == "TABLE" and self._add_table(slide, content, content_slots):
            return

        slot_text = [item if isinstance(item, str) else item.get("text", "") for item in content.get("bullets", [])] or [content.get("body", "")]
        if len(slot_text) > 1:
            compact_slots = [slot for slot in content_slots if slot["w"] < 8]
            content_slots = compact_slots or content_slots
        if len(content_slots) < len(slot_text) and content.get("body"):
            slot_text = [content["body"]]
        for item, text in zip(sorted(content_slots, key=lambda slot: (slot["y"], slot["x"])), filter(None, slot_text)):
            box = self._add_layout_textbox(slide, {"x": item["x"] + 0.18, "y": item["y"] + 0.18, "w": max(item["w"] - 0.36, 0.2), "h": max(item["h"] - 0.36, 0.2)})
            box.text_frame.paragraphs[0].text = text
            self._style_paragraph(box.text_frame.paragraphs[0], 13 if len(text) > 80 else 16, self.tokens["body_font"], bold=False)
            if self._is_dark(item["background"]):
                for run in box.text_frame.paragraphs[0].runs:
                    run.font.color.rgb = self._rgb("#FFFFFF", self.DEFAULT_COLORS["text"])

    def _template_name(self, index: int) -> str:
        config = self._as_dict(getattr(self.template_config, "config", self.template_config))
        archive = self._as_dict(config.get("zipTemplate"))
        names = archive.get("slides") if isinstance(archive.get("slides"), list) else []
        return str(names[index]).lower() if 0 <= index < len(names) else ""

    @staticmethod
    def _content_texts(content: dict) -> list[str]:
        bullets = content.get("bullets") if isinstance(content.get("bullets"), list) else []
        texts = [str(item.get("text", "")) if isinstance(item, dict) else str(item) for item in bullets]
        return [text.strip() for text in texts if text and text.strip()] or [str(content.get("body", "")).strip()]

    def _write(self, slide: Any, text: str, x: float, y: float, w: float, h: float, size: int, *, color: str = "#1A1A1A", bold: bool = False, font: Optional[str] = None) -> None:
        box = self._add_layout_textbox(slide, {"x": x, "y": y, "w": w, "h": h})
        box.text_frame.word_wrap = True
        paragraph = box.text_frame.paragraphs[0]
        paragraph.text = text
        self._style_paragraph(paragraph, size, font or self.tokens["body_font"], bold=bold)
        for run in paragraph.runs:
            run.font.color.rgb = self._rgb(color, self.DEFAULT_COLORS["text"])

    def _add_semantic_html_layout(self, slide: Any, template_name: str, slide_data: Any, content: dict) -> bool:
        """Fill imported report layouts by their information structure, not as blank decoration."""
        if not any(key in template_name for key in ("threat-model", "rsp-tier", "methodology", "external-evaluators")):
            return False
        title = str(content.get("heading") or getattr(slide_data, "title", ""))
        body = str(content.get("body") or "")
        texts = self._content_texts(content)
        self._write(slide, title, .83, .90, 10.2, .55, 30, bold=True)
        self._write(slide, body, .83, 1.48, 11.65, .35, 11, color="#5C5C5C")
        self._write(slide, "AI SECURITY REPORT", .83, .42, 4.2, .18, 8, color="#5C5C5C", bold=True, font="JetBrains Mono")
        if "threat-model" in template_name:
            self._add_threat_model(slide, texts)
        elif "rsp-tier" in template_name:
            self._add_rsp_tier(slide, texts)
        elif "methodology" in template_name:
            self._add_methodology(slide, texts)
        else:
            self._add_external_evaluators(slide, texts)
        return True

    def _add_threat_model(self, slide: Any, texts: list[str]) -> None:
        headers = ["위협 시나리오", "공격 표면", "핵심 대응", "우선순위"]
        x, y, w = .83, 2.0, 11.67
        widths = [3.0, 2.45, 4.35, 1.87]
        self._add_table_row(slide, headers, x, y, widths, .42, header=True)
        for index, text in enumerate((texts + [""] * 4)[:4]):
            parts = [part.strip() for part in text.replace("·", "-").split("-") if part.strip()]
            values = [parts[0] if parts else text, "모델·도구 접근", "정책·검증·모니터링", "높음" if index < 2 else "중간"]
            self._add_table_row(slide, values, x, y + .42 + index * .92, widths, .92, shaded=index % 2 == 1)

    def _add_table_row(self, slide: Any, values: list[str], x: float, y: float, widths: list[float], h: float, *, header: bool = False, shaded: bool = False) -> None:
        cursor = x
        for value, width in zip(values, widths):
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cursor), Inches(y), Inches(width), Inches(h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = self._rgb("#F2F1EA" if shaded else "#FAFAF7", "#FFFFFF")
            shape.line.color.rgb = self._rgb("#1A1A1A" if header else "#D5D5CD", "#D5D5CD")
            self._write(slide, value, cursor + .10, y + .10, width - .20, h - .16, 10 if header else 12, bold=header or width == widths[0])
            cursor += width

    def _add_rsp_tier(self, slide: Any, texts: list[str]) -> None:
        levels = ["ASL-1\n기본 안전 통제", "ASL-2\n위험 신호 탐지", "ASL-3\n고위험 역량 관리", "ASL-4\n상시 재평가"]
        for index, level in enumerate(levels):
            self._write(slide, level, .98 + index * 2.92, 2.13, 2.45, .55, 12, color="#FAFAF7" if index == 2 else "#1A1A1A", bold=True)
        for index in range(2):
            text = (texts + [""] * 2)[index]
            x = 1.03 + index * 5.97
            self._write(slide, f"TRIGGER 0{index + 1}", x, 3.66, 4.9, .18, 9, color="#C8541C", bold=True, font="JetBrains Mono")
            self._write(slide, text, x, 4.02, 4.9, 1.85, 16, bold=True)

    def _add_methodology(self, slide: Any, texts: list[str]) -> None:
        metrics = [("4", "핵심 통제 영역"), ("12", "검증 시나리오"), ("3", "독립 검토 단계"), ("100%", "대응 체계 점검")]
        for index, (value, label) in enumerate(metrics):
            col, row = index % 2, index // 2
            x, y = 1.02 + col * 2.25, 2.65 + row * 1.42
            self._write(slide, value, x, y, 1.9, .48, 32, color="#FAFAF7", bold=True)
            self._write(slide, label, x, y + .52, 1.95, .32, 9, color="#9C9C95")
        for index, text in enumerate((texts + [""] * 4)[:4]):
            y = 2.42 + index * 1.0
            self._write(slide, f"0{index + 1}", 6.58, y + .17, .35, .22, 10, color="#C8541C", bold=True, font="JetBrains Mono")
            self._write(slide, text, 7.0, y + .12, 5.15, .52, 13, bold=True)

    def _add_external_evaluators(self, slide: Any, texts: list[str]) -> None:
        names = ["독립 보안 검토", "레드팀 검증", "운영 감사"]
        for index, name in enumerate(names):
            x = 1.03 + index * 3.96
            text = (texts + [""] * 3)[index]
            self._write(slide, "EXTERNAL REVIEW", x, 2.20, 2.95, .18, 8, color="#5C5C5C", bold=True, font="JetBrains Mono")
            self._write(slide, name, x, 2.55, 2.95, .42, 17, bold=True)
            self._write(slide, text, x, 3.43, 2.95, 2.0, 12)
        self._write(slide, body := "독립적인 검토 결과를 운영 통제와 개선 계획에 반영합니다.", 1.0, 6.22, 11.2, .3, 11, color="#1A1A1A")

    def _add_chart(self, slide: Any, content: dict, slots: list[dict] | None = None) -> bool:
        chart = content.get("chart") if isinstance(content.get("chart"), dict) else {}
        labels, values = chart.get("labels"), chart.get("values")
        if not (isinstance(labels, list) and isinstance(values, list) and 2 <= len(labels) == len(values) <= 6):
            return False
        if not all(isinstance(label, str) and label.strip() for label in labels) or not all(isinstance(value, (int, float)) for value in values):
            return False
        data = CategoryChartData()
        data.categories = labels
        data.add_series(chart.get("series", "Value"), values)
        light_slots = [slot for slot in slots or [] if not self._is_dark(slot.get("background"))]
        slot = max(light_slots, key=lambda item: item["w"] * item["h"], default=None)
        x, y, w, h = (slot["x"] + 0.25, slot["y"] + 0.35, max(slot["w"] - 0.5, 2), max(slot["h"] - 0.7, 1.5)) if slot else (1.0, 2.0, 11.3, 4.6)
        graphic = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(x), Inches(y), Inches(w), Inches(h), data)
        graphic.chart.has_legend = False
        graphic.chart.value_axis.has_major_gridlines = True
        return True

    def _add_table(self, slide: Any, content: dict, slots: list[dict] | None = None) -> bool:
        table = content.get("table") if isinstance(content.get("table"), dict) else {}
        headers, rows = table.get("headers"), table.get("rows")
        if not (isinstance(headers, list) and 1 <= len(headers) <= 8 and all(isinstance(header, str) and header.strip() for header in headers)):
            return False
        if not (isinstance(rows, list) and 1 <= len(rows) <= 12):
            return False
        for row in rows:
            if not (isinstance(row, list) and len(row) == len(headers) and all(isinstance(cell, str) for cell in row)):
                return False
        light_slots = [slot for slot in slots or [] if not self._is_dark(slot.get("background"))]
        slot = max(light_slots, key=lambda item: item["w"] * item["h"], default=None)
        x, y, w, h = (slot["x"] + 0.25, slot["y"] + 0.35, max(slot["w"] - 0.5, 2), max(slot["h"] - 0.7, 1.5)) if slot else (1.0, 2.0, 11.3, 4.6)
        row_height = h / (len(rows) + 1)
        widths = [w / len(headers)] * len(headers)
        self._add_table_row(slide, headers, x, y, widths, row_height, header=True)
        for index, row in enumerate(rows):
            self._add_table_row(slide, row, x, y + row_height * (index + 1), widths, row_height, shaded=index % 2 == 1)
        return True

    @staticmethod
    def _is_dark(color: str | None) -> bool:
        if not color or not color.startswith("#") or len(color) != 7:
            return False
        red, green, blue = (int(color[index:index + 2], 16) for index in (1, 3, 5))
        return red * 299 + green * 587 + blue * 114 < 128000

    def _template_index(self, slide_data: Any, slide_index: int, total_slides: int) -> int:
        """Choose a matching imported layout, falling back to an even spread."""
        archive = self._as_dict(self._as_dict(getattr(self.template_config, "config", self.template_config)).get("zipTemplate"))
        names = archive.get("slides") if isinstance(archive.get("slides"), list) else []
        selected = self._as_dict(getattr(slide_data, "content", {})).get("templateIndex")
        selected_name = names[selected].lower() if isinstance(selected, int) and selected < len(names) and isinstance(names[selected], str) else ""
        title = str(getattr(slide_data, "title", "")).lower()
        # ponytail: block obvious appendix/reference choices; add semantic template metadata if names prove insufficient.
        if isinstance(selected, int) and 0 <= selected < len(self.html_slides) and (not any(word in selected_name for word in ("appendix", "reference")) or any(word in title for word in ("appendix", "reference", "참고", "부록"))):
            return selected
        keywords = {
            "TITLE": ("cover", "title", "intro"),
            "BULLET_LIST": ("agenda", "outline", "list"),
            "SECTION_HEADER": ("section", "divider", "pov"),
            "CONTENT": ("market", "strategy", "case", "roadmap", "future", "overview"),
            "QUOTE": ("executive-summary", "summary", "conclusion", "residual-risk"),
        }.get(str(getattr(slide_data, "type", "")).upper(), ())
        matches = [index for keyword in keywords for index, name in enumerate(names) if isinstance(name, str) and keyword in name.lower()]
        if matches:
            return matches[slide_index % len(matches)]
        return round(slide_index * (len(self.html_slides) - 1) / max(total_slides - 1, 1))

    def _add_title_slide(self, slide_data: Any):
        """Add title slide"""
        blank_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(blank_layout)
        self._apply_background(slide)

        content = slide_data.content
        title = content.get("heading", slide_data.title or "")
        subtitle = content.get("subheading", "")

        # Title
        title_layout = self._layout("title", {"x": 0.5, "y": 2.5, "w": 12.333, "h": 1.5, "fontSize": 54})
        title_box = self._add_layout_textbox(slide, title_layout)
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        self._style_paragraph(tf.paragraphs[0], title_layout["fontSize"], self.tokens["title_font"], bold=True)
        self._apply_alignment(tf.paragraphs[0], title_layout.get("align"))
        if "align" not in title_layout:
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        self._shrink_text_to_fit(title_box)

        # Subtitle
        if subtitle:
            subtitle_layout = self._layout("subtitle", {"x": 1, "y": 4.2, "w": 11.333, "h": 0.8, "fontSize": 24})
            sub_box = self._add_layout_textbox(slide, subtitle_layout)
            tf = sub_box.text_frame
            tf.paragraphs[0].text = subtitle
            self._style_paragraph(tf.paragraphs[0], subtitle_layout["fontSize"], self.tokens["body_font"])
            self._apply_alignment(tf.paragraphs[0], subtitle_layout.get("align"))
            if "align" not in subtitle_layout:
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            self._shrink_text_to_fit(sub_box)

    def _add_content_slide(self, slide_data: Any):
        """Add content slide with title and body"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self._apply_background(slide)

        content = slide_data.content
        title = content.get("heading", slide_data.title or "")
        body = content.get("body", "")
        bullets = content.get("bullets", [])

        # Title
        title_layout = self._layout("title", {"x": 0.5, "y": 0.3, "w": 12.333, "h": 0.8, "fontSize": 36})
        title_box = self._add_layout_textbox(slide, title_layout)
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        self._style_paragraph(tf.paragraphs[0], title_layout["fontSize"], self.tokens["title_font"], bold=True)
        self._apply_alignment(tf.paragraphs[0], title_layout.get("align"))
        self._shrink_text_to_fit(title_box)

        # ponytail: _add_table/_add_chart were only wired into the HTML-template
        # dispatch (_add_html_template_slide); a presentation with no template
        # attached — the common case — fell all the way through to this generic
        # content slide, which never looked at content["table"]/content["chart"]
        # at all, silently dropping real table/chart data down to just the title.
        slide_type = str(getattr(slide_data, "type", "")).upper()
        if slide_type == "TABLE" and self._add_table(slide, content):
            return
        if slide_type == "CHART" and self._add_chart(slide, content):
            return

        # Content area
        content_top = Inches(1.3)
        content_height = Inches(5.7)

        if body:
            body_layout = self._layout("body", {"x": 0.5, "y": 1.3, "w": 12.333, "h": 5.7, "fontSize": 20})
            body_box = self._add_layout_textbox(slide, body_layout)
            tf = body_box.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].text = body
            self._style_paragraph(tf.paragraphs[0], body_layout["fontSize"], self.tokens["body_font"])
            self._apply_alignment(tf.paragraphs[0], body_layout.get("align"))
            self._shrink_text_to_fit(body_box)

        if bullets:
            self._add_bullets(slide, bullets, content_top, content_height)

    def _add_bullet_slide(self, slide_data: Any):
        """Add slide with bullet points"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self._apply_background(slide)

        content = slide_data.content
        title = content.get("heading", slide_data.title or "")
        bullets = content.get("bullets", [])

        # Title
        title_layout = self._layout("title", {"x": 0.5, "y": 0.3, "w": 12.333, "h": 0.8, "fontSize": 36})
        title_box = self._add_layout_textbox(slide, title_layout)
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        self._style_paragraph(tf.paragraphs[0], title_layout["fontSize"], self.tokens["title_font"], bold=True)
        self._apply_alignment(tf.paragraphs[0], title_layout.get("align"))
        self._shrink_text_to_fit(title_box)

        # Bullets
        self._add_bullets(slide, bullets, Inches(1.3), Inches(5.7))

    def _add_bullets(
        self, slide, bullets: list, top: Inches, height: Inches
    ):
        """Add bullet points to a slide"""
        bullet_layout = self._layout(
            "bullets",
            {"x": 0.5, "y": top.inches, "w": 12.333, "h": height.inches, "fontSize": 20},
        )
        bullet_box = self._add_layout_textbox(slide, bullet_layout)
        tf = bullet_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if isinstance(bullet, dict):
                text = bullet.get("text", str(bullet))
                level = bullet.get("level", 0)
            else:
                text = str(bullet)
                level = 0

            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = f"• {text}"
            self._style_paragraph(p, bullet_layout["fontSize"], self.tokens["body_font"])
            self._apply_alignment(p, bullet_layout.get("align"))
            p.level = level
            p.space_before = Pt(12)

        if bullets:
            self._shrink_text_to_fit(bullet_box)

    def _add_two_column_slide(self, slide_data: Any):
        """Add two-column slide"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self._apply_background(slide)

        content = slide_data.content
        title = content.get("heading", slide_data.title or "")

        # Title
        title_layout = self._layout("title", {"x": 0.5, "y": 0.3, "w": 12.333, "h": 0.8, "fontSize": 36})
        title_box = self._add_layout_textbox(slide, title_layout)
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        self._style_paragraph(tf.paragraphs[0], title_layout["fontSize"], self.tokens["title_font"], bold=True)
        self._apply_alignment(tf.paragraphs[0], title_layout.get("align"))
        self._shrink_text_to_fit(title_box)

        rect = self._layout("columns", {"x": 0.5, "y": 1.15, "w": 12.3, "h": 5.85})
        gutter = 0.5
        column_w = (rect["w"] - gutter) / 2
        no_header_top = rect["y"] + 0.15
        bottom = rect["y"] + rect["h"]

        columns = content.get("columns")
        if isinstance(columns, list) and len(columns) == 2 and all(isinstance(item, dict) for item in columns):
            for index, column in enumerate(columns):
                x = rect["x"] + index * (column_w + gutter)
                header = str(column.get("header", "")).strip()
                bullets_top = no_header_top
                if header:
                    header_box = self._add_layout_textbox(slide, {"x": x, "y": rect["y"], "w": column_w, "h": 0.45})
                    header_paragraph = header_box.text_frame.paragraphs[0]
                    header_paragraph.text = header
                    self._style_paragraph(header_paragraph, 16, self.tokens["body_font"], bold=True)
                    bullets_top = rect["y"] + 0.55
                    self._shrink_text_to_fit(header_box)
                bullets = column.get("bullets") if isinstance(column.get("bullets"), list) else []
                self._add_column_bullets(slide, bullets, x, bullets_top, bottom - bullets_top, column_w)
            return

        # No columns: fall back to splitting a flat bullets array in half.
        bullets = content.get("bullets", [])
        mid = len(bullets) // 2
        left_bullets = bullets[:mid] if mid > 0 else bullets
        right_bullets = bullets[mid:] if mid > 0 else []
        self._add_column_bullets(slide, left_bullets, rect["x"], no_header_top, bottom - no_header_top, column_w)
        self._add_column_bullets(slide, right_bullets, rect["x"] + column_w + gutter, no_header_top, bottom - no_header_top, column_w)

    def _add_column_bullets(self, slide: Any, bullets: list, x: float, top: float, height: float, width: float = 5.9) -> None:
        if not bullets:
            return
        box = slide.shapes.add_textbox(Inches(x), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        for index, bullet in enumerate(bullets):
            if isinstance(bullet, dict):
                text = bullet.get("text", str(bullet))
                level = bullet.get("level", 0)
            else:
                text = str(bullet)
                level = 0
            paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
            paragraph.text = f"• {text}"
            self._style_paragraph(paragraph, 18, self.tokens["body_font"])
            paragraph.level = level if isinstance(level, int) else 0
            paragraph.space_before = Pt(10)
        self._shrink_text_to_fit(box)

    def _add_slide_title(self, slide: Any, slide_data: Any, content: dict) -> None:
        """Draw the standard top title textbox shared by the four new layouts."""
        title = content.get("heading", slide_data.title or "")
        title_layout = self._layout("title", {"x": 0.5, "y": 0.3, "w": 12.333, "h": 0.8, "fontSize": 36})
        title_box = self._add_layout_textbox(slide, title_layout)
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        self._style_paragraph(tf.paragraphs[0], title_layout["fontSize"], self.tokens["title_font"], bold=True)
        self._apply_alignment(tf.paragraphs[0], title_layout.get("align"))
        self._shrink_text_to_fit(title_box)

    def _add_timeline_slide(self, slide_data: Any):
        """Add a horizontal timeline/roadmap slide."""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self._apply_background(slide)

        content = slide_data.content
        self._add_slide_title(slide, slide_data, content)

        items = content["timeline"]["items"]
        count = len(items)
        rect = self._layout("timeline", {"x": 1.0, "y": 3.05, "w": 11.333, "h": 2.65})
        left, right = rect["x"], rect["x"] + rect["w"]
        line_y = rect["y"] + rect["h"] * (0.55 / 2.65)
        date_top, date_h = rect["y"], rect["h"] * (0.4 / 2.65)
        text_top, text_h = rect["y"] + rect["h"] * (0.85 / 2.65), rect["h"] * (1.8 / 2.65)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(line_y), Inches(right - left), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = self.tokens["text"]
        line.line.fill.background()

        slot_w = (right - left) / count
        for index, item in enumerate(items):
            cx = left + slot_w * index + slot_w / 2
            marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 0.09), Inches(line_y - 0.07), Inches(0.18), Inches(0.18))
            marker.fill.solid()
            marker.fill.fore_color.rgb = self.tokens["text"]
            marker.line.fill.background()

            date = str(item.get("date", "")).strip()
            if date:
                date_box = self._add_layout_textbox(slide, {"x": cx - slot_w / 2 + 0.05, "y": date_top, "w": slot_w - 0.1, "h": date_h})
                date_box.text_frame.word_wrap = True
                date_paragraph = date_box.text_frame.paragraphs[0]
                date_paragraph.text = date
                self._style_paragraph(date_paragraph, 11, self.tokens["body_font"], bold=True)
                date_paragraph.alignment = PP_ALIGN.CENTER
                self._shrink_text_to_fit(date_box)

            label = str(item.get("label", "")).strip()
            description = str(item.get("description", "")).strip()
            text_box = self._add_layout_textbox(slide, {"x": cx - slot_w / 2 + 0.05, "y": text_top, "w": slot_w - 0.1, "h": text_h})
            text_box.text_frame.word_wrap = True
            label_paragraph = text_box.text_frame.paragraphs[0]
            label_paragraph.text = label
            self._style_paragraph(label_paragraph, 13, self.tokens["body_font"], bold=True)
            label_paragraph.alignment = PP_ALIGN.CENTER
            if description:
                description_paragraph = text_box.text_frame.add_paragraph()
                description_paragraph.text = description
                self._style_paragraph(description_paragraph, 11, self.tokens["body_font"])
                description_paragraph.alignment = PP_ALIGN.CENTER
            self._shrink_text_to_fit(text_box)

    def _add_process_slide(self, slide_data: Any):
        """Add a left-to-right numbered process/step-flow slide."""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self._apply_background(slide)

        content = slide_data.content
        self._add_slide_title(slide, slide_data, content)

        steps = content["process"]["steps"]
        count = len(steps)
        rect = self._layout("process", {"x": 0.7, "y": 2.8, "w": 11.933, "h": 1.8})
        left, right, y, h, gap = rect["x"], rect["x"] + rect["w"], rect["y"], rect["h"], 0.4
        box_w = (right - left - gap * (count - 1)) / count
        for index, step in enumerate(steps):
            x = left + index * (box_w + gap)
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_w), Inches(h))
            box.fill.solid()
            box.fill.fore_color.rgb = self.tokens["background"]
            box.line.color.rgb = self.tokens["text"]
            box_tf = box.text_frame
            box_tf.word_wrap = True
            number_paragraph = box_tf.paragraphs[0]
            number_paragraph.text = f"{index + 1}. {step.get('label', '')}"
            self._style_paragraph(number_paragraph, 14, self.tokens["body_font"], bold=True)
            number_paragraph.alignment = PP_ALIGN.CENTER
            description = str(step.get("description", "")).strip()
            if description:
                description_paragraph = box_tf.add_paragraph()
                description_paragraph.text = description
                self._style_paragraph(description_paragraph, 11, self.tokens["body_font"])
                description_paragraph.alignment = PP_ALIGN.CENTER
            self._shrink_text_to_fit(box)
            if index < count - 1:
                arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x + box_w), Inches(y + h / 2 - 0.15), Inches(gap), Inches(0.3))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = self.tokens["text"]
                arrow.line.fill.background()

    def _add_comparison_slide(self, slide_data: Any):
        """Add a two-sided VS comparison slide."""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self._apply_background(slide)

        content = slide_data.content
        self._add_slide_title(slide, slide_data, content)

        rect = self._layout("comparison", {"x": 0.5, "y": 1.3, "w": 12.3, "h": 5.7})
        gutter = 0.5
        column_w = (rect["w"] - gutter) / 2
        comparison = content["comparison"]
        for side_key, x in (("left", rect["x"]), ("right", rect["x"] + column_w + gutter)):
            side = comparison[side_key]
            header_box = self._add_layout_textbox(slide, {"x": x, "y": rect["y"], "w": column_w, "h": 0.5})
            header_paragraph = header_box.text_frame.paragraphs[0]
            header_paragraph.text = str(side.get("title", ""))
            self._style_paragraph(header_paragraph, 20, self.tokens["body_font"], bold=True)
            header_paragraph.alignment = PP_ALIGN.CENTER
            bullets_top = rect["y"] + 0.6
            self._add_column_bullets(slide, side.get("bullets", []), x, bullets_top, rect["y"] + rect["h"] - bullets_top, column_w)

        badge_x = rect["x"] + column_w + gutter / 2 - 0.4
        badge_y = rect["y"] + 0.05
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(badge_x), Inches(badge_y), Inches(0.8), Inches(0.8))
        badge.fill.solid()
        badge.fill.fore_color.rgb = self.tokens["text"]
        badge.line.fill.background()
        badge_paragraph = badge.text_frame.paragraphs[0]
        badge_paragraph.text = "VS"
        self._style_paragraph(badge_paragraph, 16, self.tokens["body_font"], bold=True)
        badge_paragraph.alignment = PP_ALIGN.CENTER
        for run in badge_paragraph.runs:
            # The badge fill is tokens["text"] (dark by default) so its own text
            # needs to contrast against that fill, not the slide background.
            run.font.color.rgb = self.tokens["background"]

    def _add_kpi_slide(self, slide_data: Any):
        """Add a grid of KPI metric cards."""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self._apply_background(slide)

        content = slide_data.content
        self._add_slide_title(slide, slide_data, content)

        metrics = content["metrics"]["metrics"]
        count = len(metrics)
        columns = 3 if count > 4 else 2
        rows = math.ceil(count / columns)
        rect = self._layout("kpi", {"x": 0.7, "y": 1.6, "w": 11.933, "h": 5.3})
        left, top, right, bottom, gap = rect["x"], rect["y"], rect["x"] + rect["w"], rect["y"] + rect["h"], 0.3
        card_w = (right - left - gap * (columns - 1)) / columns
        card_h = (bottom - top - gap * (rows - 1)) / rows
        for index, metric in enumerate(metrics):
            col, row = index % columns, index // columns
            # Center a ragged last row (e.g. 5 cards over 3 columns) instead of
            # leaving its cards flush left with an empty gap on the right.
            row_count = min(columns, count - row * columns)
            row_offset = (columns - row_count) * (card_w + gap) / 2
            x = left + row_offset + col * (card_w + gap)
            y = top + row * (card_h + gap)
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_w), Inches(card_h))
            card.fill.solid()
            card.fill.fore_color.rgb = self.tokens["background"]
            card.line.color.rgb = self.tokens["text"]
            card_tf = card.text_frame
            card_tf.word_wrap = True
            value_paragraph = card_tf.paragraphs[0]
            value_paragraph.text = str(metric.get("value", ""))
            self._style_paragraph(value_paragraph, 32, self.tokens["body_font"], bold=True)
            value_paragraph.alignment = PP_ALIGN.CENTER
            label_paragraph = card_tf.add_paragraph()
            label_paragraph.text = str(metric.get("label", ""))
            self._style_paragraph(label_paragraph, 13, self.tokens["body_font"])
            label_paragraph.alignment = PP_ALIGN.CENTER
            self._shrink_text_to_fit(card)

    def _add_quote_slide(self, slide_data: Any):
        """Add quote slide"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self._apply_background(slide)

        content = slide_data.content
        quote_text = content.get("body", content.get("heading", ""))

        # Quote text
        body_layout = self._layout("body", {"x": 1.5, "y": 2.5, "w": 10.333, "h": 2, "fontSize": 32})
        quote_box = self._add_layout_textbox(slide, body_layout)
        tf = quote_box.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = f'"{quote_text}"'
        self._style_paragraph(tf.paragraphs[0], body_layout["fontSize"], self.tokens["body_font"], italic=True)
        self._apply_alignment(tf.paragraphs[0], body_layout.get("align"))
        if "align" not in body_layout:
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        self._shrink_text_to_fit(quote_box)

    def _add_section_header_slide(self, slide_data: Any):
        """Add section header slide"""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)
        self._apply_background(slide)

        content = slide_data.content
        title = content.get("heading", slide_data.title or "")

        # Large centered title
        title_layout = self._layout("title", {"x": 0.5, "y": 3, "w": 12.333, "h": 1.5, "fontSize": 48})
        title_box = self._add_layout_textbox(slide, title_layout)
        tf = title_box.text_frame
        tf.paragraphs[0].text = title
        self._style_paragraph(tf.paragraphs[0], title_layout["fontSize"], self.tokens["title_font"], bold=True)
        self._apply_alignment(tf.paragraphs[0], title_layout.get("align"))
        if "align" not in title_layout:
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        self._shrink_text_to_fit(title_box)
