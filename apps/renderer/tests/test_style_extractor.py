from io import BytesIO
import zipfile

from fastapi.testclient import TestClient
import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from apps.renderer.src.services.style_extractor import extract_template_tokens
from apps.renderer.src.main import app


KOREAN_TEXT = "한글 제목은 추출하면 안 됩니다"


def _set_east_asian_font(run, typeface):
    r_pr = run._r.get_or_add_rPr()
    east_asian = OxmlElement("a:ea")
    east_asian.set("typeface", typeface)
    r_pr.append(east_asian)


def _example_pptx():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x11, 0x22, 0x33)

    title = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    run = title.text_frame.paragraphs[0].add_run()
    run.text = KOREAN_TEXT
    run.font.name = "Fallback Font"
    run.font.size = Pt(32)
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    _set_east_asian_font(run, "Noto Sans KR")

    for left in (1, 3):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(3), Inches(1), Inches(1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x44, 0x55, 0x66)

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_extracts_only_deterministic_style_tokens_and_prefers_east_asian_font():
    tokens = extract_template_tokens(_example_pptx())

    assert tokens == {
        "colors": {"background": "#112233", "primary": "#445566"},
        "typography": {"titleFont": "Noto Sans KR", "bodyFont": "Noto Sans KR"},
        "htmlTemplate": '<div data-jaslide-slot="title" data-x="1" data-y="1" data-w="8" data-h="1" data-font-size="32" data-align="center"></div>',
    }
    assert KOREAN_TEXT not in str(tokens)


def test_ignores_patterned_shape_fills():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(1), Inches(1))
    shape.fill.patterned()
    shape.fill.fore_color.rgb = RGBColor(0xAA, 0xBB, 0xCC)

    buffer = BytesIO()
    presentation.save(buffer)

    assert extract_template_tokens(buffer.getvalue()) == {"colors": {}, "typography": {}}


def test_extract_style_upload_returns_only_config_tokens():
    response = TestClient(app).post(
        "/api/extract/style",
        files={
            "file": (
                "example.pptx",
                _example_pptx(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )

    assert response.status_code == 200
    body = response.json()

    # The endpoint returns one `config` object and nothing else — no file, no envelope.
    # Its contents grew when a PPTX upload became an editable template: alongside the
    # style tokens it carries the per-slide HTML and the object map the editor needs,
    # which is what `source.kind == "pptx"` selects on. This test used to pin the
    # pre-PPTX shape and had been failing ever since.
    assert list(body) == ["config"]
    config = body["config"]
    assert config["colors"] == {"background": "#112233", "primary": "#445566"}
    assert config["typography"] == {"titleFont": "Noto Sans KR", "bodyFont": "Noto Sans KR"}
    assert config["source"]["kind"] == "pptx"
    assert len(config["htmlSlides"]) == 1
    assert config["archive"]["canvas"] == {"width": 1920, "height": 1080}
    # The text box and both rectangles must be addressable, or the editor has nothing
    # to select on this slide.
    kinds = [obj["kind"] for obj in config["source"]["slides"][0]["objects"]]
    assert kinds.count("text") == 1 and kinds.count("shape") == 2
    assert KOREAN_TEXT in config["htmlSlides"][0], "the slide's own words must survive for editing"
    # The style tokens themselves stay free of deck content.
    assert KOREAN_TEXT not in str(config["colors"]) + str(config["typography"])


def test_extract_content_returns_text_by_slide_without_style_tokens():
    response = TestClient(app).post(
        "/api/extract/content",
        files={
            "file": (
                "example.pptx",
                _example_pptx(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )

    assert response.status_code == 200
    assert response.json() == {
        "content": KOREAN_TEXT,
        "slides": [{"number": 1, "title": KOREAN_TEXT, "content": KOREAN_TEXT}],
    }


@pytest.mark.parametrize("filename", ["example.txt", "example.pdf"])
def test_extract_style_rejects_non_pptx_name(filename):
    response = TestClient(app).post(
        "/api/extract/style",
        files={"file": (filename, _example_pptx(), "application/octet-stream")},
    )

    assert response.status_code == 400
    assert response.json() == {"detail": "PPTX file required"}


def test_extract_style_accepts_generic_browser_pptx_mime_type():
    response = TestClient(app).post(
        "/api/extract/style",
        files={"file": ("example.pptx", _example_pptx(), "application/octet-stream")},
    )

    assert response.status_code == 200


def _zip_package(entries):
    buffer = BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as package:
        for name, content in entries.items():
            package.writestr(name, content)
    return buffer.getvalue()


def test_extract_style_rejects_invalid_zip_before_parsing():
    response = TestClient(app).post(
        "/api/extract/style",
        files={
            "file": (
                "example.pptx",
                b"not a zip archive",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )

    assert response.status_code == 400
    assert response.json() == {"detail": "Invalid PPTX package"}


def test_extract_style_rejects_pptx_zip_with_expanded_content_over_limit():
    entries = {
        "[Content_Types].xml": b"<Types/>",
        "_rels/.rels": b"<Relationships/>",
        "ppt/presentation.xml": b"<p:presentation/>",
    }
    chunk = b"x" * (1024 * 1024 + 1)
    entries.update({f"ppt/media/{index}.bin": chunk for index in range(100)})
    response = TestClient(app).post(
        "/api/extract/style",
        files={
            "file": (
                "example.pptx",
                _zip_package(entries),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )

    assert response.status_code == 400
    assert response.json() == {"detail": "Invalid PPTX package"}
