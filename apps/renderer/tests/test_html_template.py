from apps.renderer.src.services.html_template import parse_html_layout
from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

from apps.renderer.src.main import app


def test_parse_html_layout_reads_a_title_slot():
    layout = parse_html_layout(
        '<h1 data-jaslide-slot="title" data-x="1" data-y="2" data-w="10" '
        'data-h="1" data-font-size="42" data-align="center"></h1>'
    )

    assert layout == {
        "title": {
            "x": 1.0,
            "y": 2.0,
            "w": 10.0,
            "h": 1.0,
            "fontSize": 42,
            "align": "center",
        }
    }


def test_parse_html_layout_ignores_unknown_slots_and_invalid_rectangles():
    assert parse_html_layout(
        '<p data-jaslide-slot="script" data-x="1" data-y="1" data-w="1" data-h="1"></p>'
    ) == {}
    assert parse_html_layout(
        '<p data-jaslide-slot="body" data-x="12" data-y="1" data-w="2" data-h="1"></p>'
    ) == {}


def test_render_request_keeps_the_html_template_layout():
    response = TestClient(app).post(
        "/api/render/pptx",
        json={
            "presentation": {
                "id": "presentation-1",
                "title": "Template check",
                "template": {
                    "id": "template-1",
                    "name": "HTML template",
                    "config": {
                        "htmlTemplate": (
                            '<h1 data-jaslide-slot="title" data-x="1" data-y="1" '
                            'data-w="9" data-h="1"></h1>'
                        )
                    },
                },
                "slides": [
                    {
                        "id": "slide-1",
                        "order": 0,
                        "type": "CONTENT",
                        "content": {"heading": "Heading"},
                    }
                ],
            }
        },
    )

    assert response.status_code == 200
    generated = Presentation(BytesIO(response.content))
    assert generated.slides[0].shapes[0].left == Inches(1)
