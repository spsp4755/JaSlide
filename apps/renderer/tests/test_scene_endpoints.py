import base64

from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO

from apps.renderer.src.main import app

client = TestClient(app)


def _pptx_bytes() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1)).text_frame.paragraphs[0].add_run().text = "x"
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_pptx_scene_load_applies_stored_edits():
    source = base64.b64encode(_pptx_bytes()).decode("ascii")
    response = client.post("/api/scene/pptx/load", json={"sourcePptx": source, "templateIndex": 0, "objectEdits": []})

    assert response.status_code == 200
    scene = response.json()["scene"]
    assert scene["objects"][0]["type"] == "text"


def test_pptx_scene_save_returns_object_edits():
    scene = {
        "width": 1920, "height": 1080,
        "objects": [{
            "id": "shape-1", "x": 0, "y": 0, "width": 100, "height": 50, "rotation": 0,
            "sourceRef": {"kind": "pptx-shape", "shapeId": "shape-1"},
            "type": "shape", "shape": "rect", "fill": "#FFFFFF", "stroke": "#000000", "strokeWidth": 1,
        }],
    }
    response = client.post("/api/scene/pptx/save", json={"scene": scene})

    assert response.status_code == 200
    edits = response.json()["objectEdits"]
    assert edits[0]["objectId"] == "shape-1"


def test_html_scene_load_and_save_round_trip():
    html = (
        '<div class="slide-container" style="position:relative;width:1920px;height:1080px">'
        '<div data-object="true" data-object-type="textbox" '
        'style="position:absolute;left:10px;top:10px;width:200px;height:60px;color:#000000">Hi</div>'
        "</div>"
    )
    load_response = client.post("/api/scene/html/load", json={"html": html})
    assert load_response.status_code == 200
    scene = load_response.json()["scene"]

    save_response = client.post("/api/scene/html/save", json={"scene": scene})
    assert save_response.status_code == 200
    assert ">Hi<" in save_response.json()["html"]
