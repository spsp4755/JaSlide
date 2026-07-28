from io import BytesIO
from unittest.mock import patch

from fastapi.testclient import TestClient
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt

from apps.renderer.src.main import app
from apps.renderer.src.services.template_fidelity import fidelity_report


def _deck_bytes(build) -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    build(slide)
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_a_recognized_autoshape_is_not_flagged():
    source = _deck_bytes(lambda slide: slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)))

    report = fidelity_report(source)

    assert report["degradedObjects"] == []


def test_a_connector_is_flagged_because_its_arrow_direction_is_not_extracted():
    source = _deck_bytes(lambda slide: slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(1), Inches(4), Inches(1)))

    report = fidelity_report(source)

    assert len(report["degradedObjects"]) == 1
    assert report["degradedObjects"][0]["type"] == "line"


def test_a_font_installed_on_this_server_is_not_reported_missing():
    def build(slide):
        run = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame.paragraphs[0].add_run()
        run.text = "Test"
        run.font.name = "NanumGothic"
    source = _deck_bytes(build)

    with patch("apps.renderer.src.services.template_fidelity._installed_font_families", return_value=frozenset({"NanumGothic"})):
        report = fidelity_report(source)

    assert report["missingFontFamilies"] == []


def test_a_font_this_server_does_not_have_is_reported_missing():
    def build(slide):
        run = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame.paragraphs[0].add_run()
        run.text = "Test"
        run.font.name = "HY헤드라인M"
    source = _deck_bytes(build)

    with patch("apps.renderer.src.services.template_fidelity._installed_font_families", return_value=frozenset({"NanumGothic"})):
        report = fidelity_report(source)

    assert report["missingFontFamilies"] == ["HY헤드라인M"]


def test_a_table_cells_font_is_checked_too():
    def build(slide):
        table = slide.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(4), Inches(1)).table
        run = table.cell(0, 0).text_frame.paragraphs[0].add_run()
        run.text = "Cell"
        run.font.name = "Missing Font"
    source = _deck_bytes(build)

    with patch("apps.renderer.src.services.template_fidelity._installed_font_families", return_value=frozenset()):
        report = fidelity_report(source)

    assert "Missing Font" in report["missingFontFamilies"]


def test_extract_fidelity_upload_returns_the_report():
    source = _deck_bytes(lambda slide: slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(1), Inches(4), Inches(1)))

    response = TestClient(app).post(
        "/api/extract/fidelity",
        files={"file": ("example.pptx", source, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
    )

    assert response.status_code == 200
    body = response.json()
    assert list(body) == ["degradedObjects", "missingFontFamilies"]
    assert body["degradedObjects"][0]["type"] == "line"


def test_extract_fidelity_rejects_non_pptx_name():
    response = TestClient(app).post(
        "/api/extract/fidelity",
        files={"file": ("example.zip", BytesIO(b"not a real file"), "application/zip")},
    )

    assert response.status_code == 400


def test_installed_font_families_reads_real_fc_list_output():
    # Not mocked here — this is the one test that exercises the real subprocess
    # call, so a change to its parsing breaks a test instead of shipping silently.
    from apps.renderer.src.services.template_fidelity import _installed_font_families
    _installed_font_families.cache_clear()

    families = _installed_font_families()

    assert isinstance(families, frozenset)
    assert len(families) > 0
