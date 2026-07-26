import base64
from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from apps.renderer.src.services.html_template import parse_html_objects
from apps.renderer.src.services.pptx_to_html import pptx_to_html

PNG = base64.b64decode("iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVQIHWP4z8DwHwAFgAI/ScL1nQAAAABJRU5ErkJggg==")


def test_converts_pptx_shapes_and_text_to_positioned_html_slides():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x11, 0x22, 0x33)
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "안전 <검증>"
    run.font.size = Pt(32)
    rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), Inches(3), Inches(3), Inches(1))
    rectangle.fill.solid()
    rectangle.fill.fore_color.rgb = RGBColor(0x44, 0x55, 0x66)
    buffer = BytesIO()
    presentation.save(buffer)

    result = pptx_to_html(buffer.getvalue())

    assert result["archive"]["slides"] == ["slide-01"]
    assert "background:#112233" in result["htmlSlides"][0]
    assert "left:192px" in result["htmlSlides"][0]
    assert "안전 &lt;검증&gt;" in result["htmlSlides"][0]
    assert "background:#445566" in result["htmlSlides"][0]


def test_preserves_pptx_font_family_in_html():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Weekly report"
    run.font.name = "NanumGothic"
    buffer = BytesIO()
    presentation.save(buffer)

    result = pptx_to_html(buffer.getvalue())

    assert "font-family:NanumGothic" in result["htmlSlides"][0]


def test_emits_font_sizes_in_canvas_px_not_points():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "제목"
    run.font.size = Pt(22)
    buffer = BytesIO()
    presentation.save(buffer)

    result = pptx_to_html(buffer.getvalue())

    # The canvas is 1080px over a 7.5in slide — 144px per inch, so a point is two px.
    # Emitting "22px" for 22pt drew every extracted deck's text at half size.
    assert "font-size:44px" in result["htmlSlides"][0]
    assert "font-size:22px" not in result["htmlSlides"][0]
    # The object map stays in points, which is what python-pptx and the editor read.
    assert result["source"]["slides"][0]["objects"][0]["fontSize"] == 22


def test_round_trips_a_font_size_through_html_back_to_its_own_points():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "제목"
    run.font.size = Pt(22)
    buffer = BytesIO()
    presentation.save(buffer)

    extracted = pptx_to_html(buffer.getvalue())
    objects = parse_html_objects(extracted["htmlSlides"][0])

    # 22pt out, 22pt back. The two conversions used to disagree (px==pt one way,
    # px*0.54 the other), so a deck came back at 54% of the size it left at.
    assert [item["fontSize"] for item in objects] == [22]


def test_round_trips_a_size_below_the_old_eighteen_point_floor():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(0.4))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "AI엔지니어링 파트"
    run.font.size = Pt(13)
    buffer = BytesIO()
    presentation.save(buffer)

    extracted = pptx_to_html(buffer.getvalue())

    # The container size seeded at 18 acted as a floor, so a 13pt caption came back
    # from the round trip two sizes larger than it went in.
    assert [item["fontSize"] for item in parse_html_objects(extracted["htmlSlides"][0])] == [13]
    assert "font-size:26px" in extracted["htmlSlides"][0]


def test_converts_tables_without_assuming_a_shape_fill():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2)).table
    table.cell(0, 0).text = "Header"
    table.cell(1, 1).text = "Value"
    buffer = BytesIO()
    presentation.save(buffer)

    result = pptx_to_html(buffer.getvalue())

    assert 'data-object-type="table"' in result["htmlSlides"][0]
    assert "Header" in result["htmlSlides"][0]
    assert result["source"]["kind"] == "pptx"
    assert result["source"]["slides"][0]["objects"][0]["kind"] == "table"
    assert result["source"]["slides"][0]["objects"][0]["id"]
    assert result["source"]["slides"][0]["objects"][0]["left"] == 192
    assert result["source"]["slides"][0]["objects"][0]["cells"] == [["Header", ""], ["", "Value"]]


def test_preserves_table_cell_dimensions_and_formatting():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2)).table
    table.columns[0].width = Inches(3)
    table.cell(0, 0).fill.solid()
    table.cell(0, 0).fill.fore_color.rgb = RGBColor(0x11, 0x22, 0x33)
    run = table.cell(0, 0).text_frame.paragraphs[0].add_run()
    run.text = "Header"
    run.font.name = "NanumGothic"
    run.font.size = Pt(18)
    buffer = BytesIO()
    presentation.save(buffer)

    result = pptx_to_html(buffer.getvalue())

    # Widening column 0 to 3in makes the table 5in wide; python-pptx leaves the
    # graphicFrame's stored extent at the original 4in. PowerPoint draws the table at
    # its own 5in, so the column is 3/5 of it — measuring against the stale frame
    # (3/4 = 75%) put the preview and the editor's cell grid out of step with the deck.
    assert "width:60.0%" in result["htmlSlides"][0]
    table = next(obj for obj in result["source"]["slides"][0]["objects"] if obj["kind"] == "table")
    assert table["width"] == sum(table["columnWidths"])
    assert table["height"] == sum(table["rowHeights"])
    assert "background:#112233" in result["htmlSlides"][0]
    assert "font-family:NanumGothic" in result["htmlSlides"][0]


def test_extracts_real_row_and_column_geometry_for_a_table():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(6), Inches(3)).table
    table.rows[0].height = Inches(1)
    table.rows[1].height = Inches(2)
    table.columns[0].width = Inches(4)
    table.columns[1].width = Inches(2)

    buffer = BytesIO()
    presentation.save(buffer)

    result = pptx_to_html(buffer.getvalue())
    table_object = next(obj for obj in result["source"]["slides"][0]["objects"] if obj["kind"] == "table")

    assert table_object["rowHeights"][0] < table_object["rowHeights"][1]
    assert table_object["columnWidths"][0] > table_object["columnWidths"][1]
    assert round(table_object["rowHeights"][0] / sum(table_object["rowHeights"]), 2) == round(1 / 3, 2)


def test_extracts_alignment_for_a_text_object():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text = "Heading"
    box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    buffer = BytesIO()
    presentation.save(buffer)

    result = pptx_to_html(buffer.getvalue())
    text_object = next(obj for obj in result["source"]["slides"][0]["objects"] if obj["kind"] == "text")

    assert text_object["align"] == "center"


def _bulleted(paragraph, char: str, mar_l: int, indent: int) -> None:
    """Give `paragraph` the buChar/marL/indent a real deck's list line carries."""
    properties = paragraph._pPr if paragraph._pPr is not None else paragraph._p.get_or_add_pPr()
    properties.set("marL", str(mar_l))
    properties.set("indent", str(indent))
    bullet = properties.makeelement(qn("a:buChar"), {"char": char})
    properties.append(bullet)


def test_renders_list_paragraphs_with_their_marker_and_indent():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    frame = box.text_frame
    frame.paragraphs[0].text = "IT 운영 및 AI 연구"
    _bulleted(frame.paragraphs[0], "•", 88900, 0)
    second = frame.add_paragraph()
    second.text = "프로젝트 관리 및 지원"
    second.level = 1
    _bulleted(second, "-", 481013, -285750)
    buffer = BytesIO()
    presentation.save(buffer)

    html = pptx_to_html(buffer.getvalue())["htmlSlides"][0]

    # The deck states its own bullets and indents; joining paragraphs with <br> threw
    # all of it away and left a flat block of unmarked lines.
    assert ">•<" in html and ">-<" in html
    # 88900 EMU is 14px on this canvas, 481013 is 76px, and the hanging indent is -45px.
    assert "margin-left:14px;text-indent:0px" in html
    assert "margin-left:76px;text-indent:-45px" in html


def test_keeps_a_hanging_indent_inside_the_shape():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    box.text_frame.paragraphs[0].text = "주간업무 추진실적"
    # A deck can hang a bullet further left than the margin it hangs from. CSS pulled
    # the first line outside the box, where overflow:hidden ate the opening character.
    _bulleted(box.text_frame.paragraphs[0], "•", 0, -177800)
    buffer = BytesIO()
    presentation.save(buffer)

    html = pptx_to_html(buffer.getvalue())["htmlSlides"][0]

    assert "margin-left:0px;text-indent:0px" in html
    assert "text-indent:-" not in html


def test_draws_no_marker_on_an_empty_list_paragraph():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    frame = box.text_frame
    frame.paragraphs[0].text = "실적"
    _bulleted(frame.paragraphs[0], "-", 481013, -285750)
    blank = frame.add_paragraph()
    _bulleted(blank, "-", 481013, -285750)
    buffer = BytesIO()
    presentation.save(buffer)

    html = pptx_to_html(buffer.getvalue())["htmlSlides"][0]

    # Spacing paragraphs are common at the end of a cell; PowerPoint draws no bullet
    # on one, and a marker alone on a blank line reads as a mistake.
    assert html.count(">-<") == 1


def test_extracts_a_picture_placeholder_as_an_image_object():
    # Real decks drop screenshots into a layout's picture placeholder, whose
    # shape_type reports PLACEHOLDER rather than PICTURE. Matching only on
    # PICTURE degraded those slides to an empty white box, losing the image.
    presentation = Presentation()
    layout = next(
        item for item in presentation.slide_layouts
        if any(place.placeholder_format.type == PP_PLACEHOLDER.PICTURE for place in item.placeholders)
    )
    slide = presentation.slides.add_slide(layout)
    placeholder = next(item for item in slide.placeholders if item.placeholder_format.type == PP_PLACEHOLDER.PICTURE)
    placeholder.insert_picture(BytesIO(PNG))

    buffer = BytesIO()
    presentation.save(buffer)

    result = pptx_to_html(buffer.getvalue())

    assert "image" in [obj["kind"] for obj in result["source"]["slides"][0]["objects"]]
    assert 'data-object-type="image"' in result["htmlSlides"][0]


def test_defaults_alignment_to_left_when_unset():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text = "Body"

    buffer = BytesIO()
    presentation.save(buffer)

    result = pptx_to_html(buffer.getvalue())
    text_object = next(obj for obj in result["source"]["slides"][0]["objects"] if obj["kind"] == "text")

    assert text_object["align"] == "left"


def test_a_grouped_diagram_becomes_editable_objects_not_one_opaque_box():
    # Real decks group things constantly, and a GroupShape has no `line` at all —
    # reading it unguarded raised AttributeError, which the API reported as
    # "Invalid PPTX file". Most real decks were rejected outright because of it.
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    box.text_frame.paragraphs[0].text = "그룹 안의 글자"
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(1), Inches(2), Inches(1))
    slide.shapes._spTree.remove(box._element)
    slide.shapes._spTree.remove(shape._element)
    group = slide.shapes.add_group_shape()
    group._element.append(box._element)
    group._element.append(shape._element)

    buffer = BytesIO()
    presentation.save(buffer)

    result = pptx_to_html(buffer.getvalue())

    objects = result["source"]["slides"][0]["objects"]
    assert [obj["kind"] for obj in objects] == ["text", "shape"], "the group's members, not the group"
    assert objects[0]["text"] == "그룹 안의 글자"


def test_one_unreadable_shape_does_not_cost_the_whole_upload():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))
    good = slide.shapes.add_textbox(Inches(4), Inches(1), Inches(3), Inches(1))
    good.text_frame.paragraphs[0].text = "살아남아야 함"
    # Strip the first shape's geometry so reading its position raises.
    slide.shapes[0]._element.spPr.remove(slide.shapes[0]._element.spPr.xfrm)

    buffer = BytesIO()
    presentation.save(buffer)

    result = pptx_to_html(buffer.getvalue())

    texts = [obj.get("text") for obj in result["source"]["slides"][0]["objects"]]
    assert "살아남아야 함" in texts


def test_a_large_image_is_not_inlined_as_base64():
    # A photo-heavy deck inlined every blob and produced a 40MB config, which then has
    # to live in a database row and load in the editor.
    from apps.renderer.src.services.pptx_to_html import MAX_INLINE_IMAGE_BYTES

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    big = _png_of_at_least(MAX_INLINE_IMAGE_BYTES + 1)
    slide.shapes.add_picture(BytesIO(big), Inches(1), Inches(1), Inches(4), Inches(3))

    buffer = BytesIO()
    presentation.save(buffer)

    result = pptx_to_html(buffer.getvalue())

    html = result["htmlSlides"][0]
    assert 'data-object-type="image"' in html, "the image still occupies its place"
    assert "base64" not in html
    assert len(html) < len(big) / 4


def _png_of_at_least(size: int) -> bytes:
    """A valid PNG padded past `size` with an ancillary chunk python-pptx accepts."""
    import struct
    import zlib

    header = b"\x89PNG\r\n\x1a\n"
    def chunk(kind: bytes, payload: bytes) -> bytes:
        return struct.pack(">I", len(payload)) + kind + payload + struct.pack(">I", zlib.crc32(kind + payload))

    ihdr = chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0))
    idat = chunk(b"IDAT", zlib.compress(b"\x00\x00\x00\x00"))
    padding = chunk(b"teXt", b"pad\x00" + b"x" * size)
    return header + ihdr + padding + idat + chunk(b"IEND", b"")
