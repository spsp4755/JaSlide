from io import BytesIO

from pptx import Presentation
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Inches

from apps.renderer.src.services.pptx_to_html import pptx_to_html
from apps.renderer.src.services.theme_colors import apply_transforms, theme_palette


def _deck_with_theme_filled_table() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(6), Inches(3)).table
    header = table.cell(0, 0)
    header.text = "추진실적"
    header.fill.solid()
    header.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_1
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_the_theme_palette_is_read_off_the_deck():
    palette = theme_palette(Presentation(BytesIO(_deck_with_theme_filled_table())))

    assert palette, "a deck always carries a theme"
    for slot in ("dk1", "lt1", "accent1"):
        assert slot in palette
        assert palette[slot].startswith("#") and len(palette[slot]) == 7


def test_a_theme_filled_cell_keeps_its_colour_in_the_html():
    # A real deck names a theme slot rather than an RGB value, and python-pptx
    # raises on `.rgb` for one. Reading colours that way dropped every themed
    # fill: the reference deck's grey header row came out transparent.
    html = pptx_to_html(_deck_with_theme_filled_table())["htmlSlides"][0]

    table_html = html.split("<table")[1]
    assert "background:#" in table_html, "the themed header fill was dropped"


def test_a_theme_shape_fill_survives_too():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.theme_color = MSO_THEME_COLOR.ACCENT_2
    buffer = BytesIO()
    presentation.save(buffer)

    html = pptx_to_html(buffer.getvalue())["htmlSlides"][0]

    assert 'data-object-type="shape"' in html
    shape_html = html.split('data-object-type="shape"')[1]
    assert "background:#" in shape_html.split("</div>")[0]


def test_luminance_and_tint_adjustments_are_applied():
    # "Background, at 85% luminance" is how a deck states a light grey header.
    # Ignoring the adjustment would paint it pure white and lose the banding.
    element = Presentation().slide_masters[0]._element.makeelement(qn("a:schemeClr"), {})
    lum = element.makeelement(qn("a:lumMod"), {"val": "50000"})
    element.append(lum)

    darkened = apply_transforms("#FFFFFF", element)

    assert darkened != "#FFFFFF"
    assert darkened.startswith("#")
    # Half luminance of white is mid grey, not black.
    channel = int(darkened[1:3], 16)
    assert 100 < channel < 160, darkened


def test_a_plain_rgb_colour_is_untouched_without_adjustments():
    element = Presentation().slide_masters[0]._element.makeelement(qn("a:srgbClr"), {"val": "123456"})

    assert apply_transforms("#123456", element) == "#123456"
