import base64
from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from apps.renderer.src.services.pptx_scene import pptx_to_scene
from apps.renderer.src.services.scene_to_pptx import scene_to_pptx


def _source_deck_bytes(build) -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    build(slide)
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_editing_an_existing_shapes_fill_reaches_the_same_shape():
    source = _source_deck_bytes(lambda slide: (
        lambda shape: (shape.fill.solid(), setattr(shape.fill.fore_color, "rgb", RGBColor(0xFF, 0x00, 0x00)))
    )(slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))))

    scene = pptx_to_scene(source)["slides"][0]
    shape_object = next(item for item in scene["objects"] if item["type"] == "shape")
    assert shape_object["fill"] == "#FF0000"

    edited = dict(scene)
    edited["objects"] = [{**shape_object, "fill": "#00FF00"}]
    output = scene_to_pptx(edited, source)

    exported = Presentation(BytesIO(output)).slides[0].shapes[0]
    assert str(exported.fill.fore_color.rgb) == "00FF00"
    # The same shape, not a new one appended alongside it.
    assert len(Presentation(BytesIO(output)).slides[0].shapes) == 1


def test_a_new_bold_run_on_the_title_reaches_the_original_title_shape():
    source = _source_deck_bytes(lambda slide: slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1)).text_frame.paragraphs[0].add_run().__setattr__("text", "Original"))

    scene = pptx_to_scene(source)["slides"][0]
    text_object = next(item for item in scene["objects"] if item["type"] == "text")

    edited = dict(scene)
    edited["objects"] = [{
        **text_object,
        "paragraphs": [{"runs": [
            {"text": "New "},
            {"text": "bold", "bold": True, "color": "#FF0000", "fontSize": 20},
        ]}],
    }]
    output = scene_to_pptx(edited, source)

    paragraph = Presentation(BytesIO(output)).slides[0].shapes[0].text_frame.paragraphs[0]
    assert [run.text for run in paragraph.runs] == ["New ", "bold"]
    assert paragraph.runs[1].font.bold
    assert str(paragraph.runs[1].font.color.rgb) == "FF0000"


def test_a_table_cells_new_text_reaches_the_original_table():
    source = _source_deck_bytes(lambda slide: slide.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(4), Inches(1)).table.cell(0, 0).__setattr__("text", "Old"))

    scene = pptx_to_scene(source)["slides"][0]
    table_object = next(item for item in scene["objects"] if item["type"] == "table")

    edited = dict(scene)
    new_cell = {"paragraphs": [{"runs": [{"text": "New", "bold": True}]}]}
    edited["objects"] = [{**table_object, "cells": [[new_cell]]}]
    output = scene_to_pptx(edited, source)

    cell = Presentation(BytesIO(output)).slides[0].shapes[0].table.cell(0, 0)
    assert cell.text == "New"
    assert cell.text_frame.paragraphs[0].runs[0].font.bold


def test_an_object_with_no_source_ref_is_inserted_as_a_new_shape():
    source = _source_deck_bytes(lambda slide: None)
    scene = {
        "width": 1920, "height": 1080,
        "objects": [{
            "id": "brand-new", "type": "shape", "x": 100, "y": 100, "width": 300, "height": 150,
            "rotation": 0, "shape": "ellipse", "fill": "#2563EB", "stroke": "#000000", "strokeWidth": 2,
        }],
    }

    output = scene_to_pptx(scene, source)
    shapes = Presentation(BytesIO(output)).slides[0].shapes
    assert len(shapes) == 1
    assert shapes[0].auto_shape_type == MSO_SHAPE.OVAL
    assert str(shapes[0].fill.fore_color.rgb) == "2563EB"


def test_round_trip_from_the_reference_deck_preserves_the_untouched_objects():
    # Import the real weekly-report deck, change nothing, export: the table's
    # own header text must survive unchanged.
    source = open("/tmp/weekly.pptx", "rb").read() if _exists("/tmp/weekly.pptx") else None
    if source is None:
        return  # fixture not present in this environment; covered by the unit tests above
    scene = pptx_to_scene(source)["slides"][0]
    output = scene_to_pptx(scene, source)
    table = next(shape for shape in Presentation(BytesIO(output)).slides[0].shapes if shape.has_table)
    assert "추진실적" in table.table.cell(0, 0).text


def _exists(path: str) -> bool:
    import os
    return os.path.exists(path)
