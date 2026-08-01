import base64
from io import BytesIO

import pytest
from types import SimpleNamespace

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from apps.renderer.src.generators.pptx_generator import PPTXGenerator, fit_font_scale


def _presentation(*slides):
    return SimpleNamespace(slides=slides)


def _slide(slide_type, title, content):
    return SimpleNamespace(type=slide_type, title=title, content=content)


def _runs(slide):
    return [run for shape in slide.shapes if shape.has_text_frame for paragraph in shape.text_frame.paragraphs for run in paragraph.runs]


def _rgb(color):
    return str(color.rgb)


def test_default_tokens_style_korean_title_content_and_quote_slides():
    pptx = PPTXGenerator().generate(
        _presentation(
            _slide("TITLE", "제목", {"heading": "제목", "subheading": "부제목"}),
            _slide("CONTENT", "내용", {"heading": "내용", "body": "한글 본문", "bullets": ["첫 항목"]}),
            _slide("QUOTE", "", {"body": "인용문"}),
        )
    )

    generated = Presentation(BytesIO(pptx))
    for slide in generated.slides:
        assert _rgb(slide.background.fill.fore_color) == "FFFFFF"
        for run in _runs(slide):
            assert run.font.name == "Noto Sans KR"
            assert _rgb(run.font.color) == "1E293B"


def test_template_colors_apply_to_all_slide_text_and_background_with_invalid_values_falling_back():
    template = SimpleNamespace(
        config=SimpleNamespace(
            colors={"background": "#123456", "text": "#ABCDEF", "primary": "not-a-color"},
            backgrounds={"type": "solid", "value": "#BAD"},
        )
    )
    pptx = PPTXGenerator(template).generate(
        _presentation(
            _slide("TITLE", "제목", {"heading": "제목", "subheading": "부제목"}),
            _slide("CONTENT", "내용", {"heading": "내용", "body": "한글 본문", "bullets": ["첫 항목"]}),
            _slide("QUOTE", "", {"body": "인용문"}),
        )
    )

    generated = Presentation(BytesIO(pptx))
    for slide in generated.slides:
        assert _rgb(slide.background.fill.fore_color) == "123456"
        for run in _runs(slide):
            assert _rgb(run.font.color) == "ABCDEF"


def test_two_column_slide_uses_the_shared_background_and_text_tokens():
    template = SimpleNamespace(config=SimpleNamespace(colors={"background": "#102030", "text": "#DDEEFF"}))
    pptx = PPTXGenerator(template).generate(
        _presentation(
            _slide("TWO_COLUMN", "비교", {"heading": "비교", "bullets": ["왼쪽", "오른쪽"]}),
        )
    )

    slide = Presentation(BytesIO(pptx)).slides[0]
    assert _rgb(slide.background.fill.fore_color) == "102030"
    assert {_rgb(run.font.color) for run in _runs(slide)} == {"DDEEFF"}


def test_korean_runs_use_the_east_asian_font_slot():
    pptx = PPTXGenerator().generate(
        _presentation(_slide("CONTENT", "한글 제목", {"heading": "한글 제목", "body": "한글 본문"}))
    )

    slide = Presentation(BytesIO(pptx)).slides[0]
    for run in _runs(slide):
        assert run._r.rPr.find(qn("a:ea")).get("typeface") == "Noto Sans KR"


def test_bare_template_hex_values_fall_back_to_default_tokens():
    template = SimpleNamespace(config=SimpleNamespace(colors={"background": "123456", "text": "ABCDEF"}))
    pptx = PPTXGenerator(template).generate(
        _presentation(_slide("CONTENT", "제목", {"heading": "제목", "body": "본문"}))
    )

    slide = Presentation(BytesIO(pptx)).slides[0]
    assert _rgb(slide.background.fill.fore_color) == "FFFFFF"
    assert {_rgb(run.font.color) for run in _runs(slide)} == {"1E293B"}


def test_bullets_preserve_their_text_without_prefix_artifacts():
    pptx = PPTXGenerator().generate(
        _presentation(
            _slide(
                "BULLET_LIST",
                "Agenda",
                {"heading": "Agenda", "bullets": [{"text": "First item", "level": 0}]},
            )
        )
    )

    slide = Presentation(BytesIO(pptx)).slides[0]
    paragraphs = [
        paragraph.text
        for shape in slide.shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
    ]

    assert "\u2022 First item" in paragraphs


def test_html_template_layout_positions_content_textboxes():
    template = SimpleNamespace(
        config=SimpleNamespace(
            htmlTemplate=(
                '<h1 data-jaslide-slot="title" data-x="1" data-y="1" data-w="9" '
                'data-h="1" data-align="center"></h1>'
                '<p data-jaslide-slot="body" data-x="1" data-y="2" data-w="9" '
                'data-h="3" data-font-size="22"></p>'
            )
        )
    )
    pptx = PPTXGenerator(template).generate(
        _presentation(_slide("CONTENT", "", {"heading": "Heading", "body": "Body"}))
    )

    shapes = Presentation(BytesIO(pptx)).slides[0].shapes
    assert shapes[0].left == Inches(1)
    assert shapes[0].top == Inches(1)
    assert shapes[0].text_frame.paragraphs[0].alignment == PP_ALIGN.CENTER
    assert shapes[1].top == Inches(2)
    assert shapes[1].text_frame.paragraphs[0].runs[0].font.size == Pt(22)


def test_html_template_layout_styles_bullet_slide_titles():
    template = SimpleNamespace(
        config=SimpleNamespace(
            htmlTemplate=(
                '<h1 data-jaslide-slot="title" data-x="2" data-y="1" data-w="8" '
                'data-h="1" data-font-size="26" data-align="right"></h1>'
            )
        )
    )
    pptx = PPTXGenerator(template).generate(
        _presentation(_slide("BULLET_LIST", "Agenda", {"heading": "Agenda", "bullets": ["One"]}))
    )

    title = Presentation(BytesIO(pptx)).slides[0].shapes[0]
    paragraph = title.text_frame.paragraphs[0]
    assert title.left == Inches(2)
    assert title.top == Inches(1)
    assert paragraph.alignment == PP_ALIGN.RIGHT
    assert paragraph.runs[0].font.size == Pt(26)


def test_genspark_style_html_applies_colors_fonts_and_title_position_without_slots():
    template = SimpleNamespace(
        config=SimpleNamespace(
            htmlTemplate=(
                '<body style="background:#fdfcf8;color:#0a1530;font-family:Inter">'
                '<div data-object-type="textbox" style="position:absolute;left:200px;top:230px;width:1400px;'
                'font-family:Newsreader;font-size:112px;color:#0a1530"></div></body>'
            )
        )
    )
    pptx = PPTXGenerator(template).generate(
        _presentation(_slide("CONTENT", "Heading", {"heading": "Heading", "body": "Body"}))
    )

    slide = Presentation(BytesIO(pptx)).slides[0]
    assert _rgb(slide.background.fill.fore_color) == "FDFCF8"
    assert slide.shapes[0].left == Inches(200 / 1920 * 13.333)
    assert slide.shapes[0].top == Inches(230 / 1080 * 7.5)
    assert slide.shapes[0].text_frame.paragraphs[0].runs[0].font.name == "Newsreader"


def test_html_slide_template_renders_its_shape_and_title_layout():
    template = SimpleNamespace(
        config=SimpleNamespace(
            htmlSlides=[
                '<div data-object="true" data-object-type="shape" style="position:absolute;left:0;top:0;width:1920px;height:1080px;background:#123456"></div>'
                '<div data-object="true" data-object-type="textbox" style="position:absolute;left:120px;top:172px;width:1680px;height:200px;font-size:56px;color:#FFFFFF">Template title</div>'
            ]
        )
    )
    pptx = PPTXGenerator(template).generate(
        _presentation(_slide("CONTENT", "Generated title", {"heading": "Generated title"}))
    )

    slide = Presentation(BytesIO(pptx)).slides[0]
    assert _rgb(slide.background.fill.fore_color) == "123456"
    assert slide.shapes[0].left == Inches(120 / 1920 * 13.333)
    assert slide.shapes[0].text == "Generated title"


def test_html_content_embeds_one_full_slide_browser_image(monkeypatch):
    png = base64.b64decode("iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAusB9Y9Jp5cAAAAASUVORK5CYII=")
    monkeypatch.setattr("apps.renderer.src.generators.pptx_generator.render_slide_png", lambda html: png)

    output = PPTXGenerator().generate(_presentation(_slide(
        "CONTENT", "HTML", {"html": "<main data-object=\"true\">HTML</main>"},
    )))

    slide = Presentation(BytesIO(output)).slides[0]
    picture = next(shape for shape in slide.shapes if shape.shape_type == 13)
    assert picture.left == 0 and picture.top == 0
    assert picture.width == Inches(13.333) and picture.height == Inches(7.5)


def test_editable_export_writes_real_shapes_instead_of_a_slide_picture(monkeypatch):
    # The default screenshot matches the HTML exactly, but exports one flat picture
    # per slide: the recipient cannot revise a single word. `editable` asks for the
    # same content as shapes and text.
    monkeypatch.setattr(
        "apps.renderer.src.generators.pptx_generator.render_slide_png",
        lambda html: pytest.fail("editable export must not screenshot the slide"),
    )
    template = SimpleNamespace(config={"htmlSlides": [
        '<div class="slide-container" style="width:1920px;height:1080px">'
        '<div data-object="true" data-object-type="shape" style="position:absolute;left:120px;top:300px;width:800px;height:400px;background:#EEF1FB"></div>'
        '<div data-object="true" data-object-type="textbox" style="position:absolute;left:120px;top:120px;width:1400px;height:90px;font-size:56px;color:#1A1A1A">제목</div>'
        '</div>',
    ]})

    slide_html = (
        '<div class="slide-container" style="width:1920px;height:1080px">'
        '<div data-object="true" data-object-type="textbox" style="position:absolute;left:120px;top:120px;width:1400px;height:90px;font-size:56px;color:#1A1A1A">잔여 위험 요약</div>'
        '<div data-object="true" data-object-type="shape" style="position:absolute;left:120px;top:300px;width:800px;height:400px;background:#1A1A1A"></div>'
        '<div data-object="true" data-object-type="textbox" style="position:absolute;left:160px;top:340px;width:720px;height:120px;font-size:32px">ASR 2.4%로 감소</div>'
        '</div>'
    )

    output = PPTXGenerator(template, editable=True).generate(_presentation(_slide(
        "CONTENT", "보고", {"html": slide_html, "heading": "주간 보고"},
    )))

    shapes = Presentation(BytesIO(output)).slides[0].shapes
    assert all(shape.shape_type != 13 for shape in shapes), "no slide-sized picture"
    # The slide's own words must survive, not just the template's empty layout.
    texts = [shape.text_frame.text for shape in shapes if shape.has_text_frame and shape.text_frame.text]
    assert "잔여 위험 요약" in texts
    assert "ASR 2.4%로 감소" in texts
    # Text on a dark panel must not come out black on black.
    dark_panel_text = next(shape for shape in shapes if shape.has_text_frame and shape.text_frame.text == "ASR 2.4%로 감소")
    assert str(dark_panel_text.text_frame.paragraphs[0].runs[0].font.color.rgb) == "FFFFFF"


def test_editable_export_rebuilds_an_html_table_as_a_real_table():
    # Flattening a grid into one string turned the report's ASR table into a run-on
    # line of headers with every row gone.
    slide_html = (
        '<div class="slide-container" style="width:1920px;height:1080px">'
        '<div data-object="true" data-object-type="table" style="position:absolute;left:120px;top:300px;width:1600px;height:400px">'
        '<table><thead><tr><th>공격 카테고리</th><th>시도</th><th>ASR</th></tr></thead>'
        '<tbody><tr><td>프롬프트 주입</td><td>1,200</td><td>2.4%</td></tr>'
        '<tr><td>탈옥</td><td>900</td><td>1.1%</td></tr></tbody></table></div>'
        '</div>'
    )

    output = PPTXGenerator(SimpleNamespace(config={}), editable=True).generate(
        _presentation(_slide("CONTENT", "ASR", {"html": slide_html})),
    )

    shape = next(item for item in Presentation(BytesIO(output)).slides[0].shapes if item.has_table)
    assert len(shape.table.rows) == 3 and len(shape.table.columns) == 3
    assert [cell.text for cell in shape.table.rows[0].cells] == ["공격 카테고리", "시도", "ASR"]
    assert [cell.text for cell in shape.table.rows[2].cells] == ["탈옥", "900", "1.1%"]
    assert shape.table.rows[0].cells[0].text_frame.paragraphs[0].runs[0].font.bold


def test_an_object_containing_nested_divs_does_not_end_early():
    # Arbitrary HTML nests divs everywhere. Closing on the first </div> ended the
    # object early and dropped everything after it on the slide.
    slide_html = (
        '<div class="slide-container" style="width:1920px;height:1080px">'
        '<div data-object="true" data-object-type="textbox" style="position:absolute;left:100px;top:100px;width:800px;height:200px">'
        '<div><span>앞</span></div><div>뒤</div></div>'
        '<div data-object="true" data-object-type="textbox" style="position:absolute;left:100px;top:400px;width:800px;height:200px">두 번째 객체</div>'
        '</div>'
    )

    output = PPTXGenerator(SimpleNamespace(config={}), editable=True).generate(
        _presentation(_slide("CONTENT", "중첩", {"html": slide_html})),
    )

    texts = [shape.text_frame.text for shape in Presentation(BytesIO(output)).slides[0].shapes if shape.has_text_frame]
    assert "앞 뒤" in texts, "both children belong to the one object"
    assert "두 번째 객체" in texts, "the object after it must still be found"


def test_html_template_chooses_layouts_by_slide_type_not_first_n_slides():
    def layout(color):
        return f'<div data-object="true" data-object-type="shape" style="position:absolute;left:0;top:0;width:1920px;height:1080px;background:{color}"></div>'

    template = SimpleNamespace(config=SimpleNamespace(
        htmlSlides=[layout("#111111"), layout("#222222"), layout("#333333")],
        zipTemplate={"slides": ["slides/cover.html", "slides/agenda.html", "slides/market.html"]},
    ))
    output = PPTXGenerator(template).generate(_presentation(
        _slide("TITLE", "Title", {"heading": "Title"}),
        _slide("BULLET_LIST", "Agenda", {"heading": "Agenda"}),
        _slide("CONTENT", "Market", {"heading": "Market"}),
    ))

    assert [_rgb(slide.background.fill.fore_color) for slide in Presentation(BytesIO(output)).slides] == ["111111", "222222", "333333"]


def test_html_template_uses_the_outline_selected_layout_when_present():
    def layout(color):
        return f'<div data-object="true" data-object-type="shape" style="position:absolute;left:0;top:0;width:1920px;height:1080px;background:{color}"></div>'

    template = SimpleNamespace(config=SimpleNamespace(htmlSlides=[layout("#111111"), layout("#222222"), layout("#333333")]))
    output = PPTXGenerator(template).generate(_presentation(_slide(
        "CONTENT", "Risk findings", {"heading": "Risk findings", "templateIndex": 2},
    )))

    assert _rgb(Presentation(BytesIO(output)).slides[0].background.fill.fore_color) == "333333"


def test_html_template_rejects_a_reference_layout_for_a_conclusion_slide():
    def layout(color):
        return f'<div data-object="true" data-object-type="shape" style="position:absolute;left:0;top:0;width:1920px;height:1080px;background:{color}"></div>'

    template = SimpleNamespace(config=SimpleNamespace(
        htmlSlides=[layout("#111111"), layout("#222222"), layout("#333333")],
        zipTemplate={"slides": ["slides/cover.html", "slides/executive-summary.html", "slides/references.html"]},
    ))
    output = PPTXGenerator(template).generate(_presentation(_slide(
        "QUOTE", "Conclusion", {"heading": "Conclusion", "templateIndex": 2},
    )))

    assert _rgb(Presentation(BytesIO(output)).slides[0].background.fill.fore_color) == "222222"


def test_html_template_fills_blank_card_shapes_with_generated_bullets():
    template = SimpleNamespace(config=SimpleNamespace(htmlSlides=[
        '<div data-object="true" data-object-type="shape" style="position:absolute;left:0;top:0;width:1920px;height:1080px;background:#FFFFFF"></div>'
        '<div data-object="true" data-object-type="shape" style="position:absolute;left:120px;top:300px;width:700px;height:400px;background:#F2F5FA"></div>'
        '<div data-object="true" data-object-type="shape" style="position:absolute;left:920px;top:300px;width:700px;height:400px;background:#F2F5FA"></div>'
    ]))
    output = PPTXGenerator(template).generate(_presentation(_slide(
        "BULLET_LIST", "Title", {"heading": "Title", "bullets": ["First", "Second"]},
    )))

    text = " ".join(shape.text for shape in Presentation(BytesIO(output)).slides[0].shapes if shape.has_text_frame)
    assert "First" in text and "Second" in text


def test_report_template_uses_a_table_layout_for_a_threat_model_slide():
    html = (
        '<div data-object="true" data-object-type="shape" style="position:absolute;left:0;top:0;width:1920px;height:1080px;background:#FAFAF7"></div>'
        '<div data-object="true" data-object-type="textbox" style="position:absolute;left:120px;top:130px;width:1400px;height:80px;font-size:56px;color:#1A1A1A">Template title</div>'
    )
    template = SimpleNamespace(config=SimpleNamespace(
        htmlSlides=[html], zipTemplate={"slides": ["slides/03-threat-model.html"]},
    ))
    output = PPTXGenerator(template).generate(_presentation(_slide(
        "CONTENT", "AI 위협 모델", {"heading": "AI 위협 모델", "body": "우선순위 기반 대응", "bullets": ["프롬프트 인젝션", "권한 오용"]},
    )))

    slide = Presentation(BytesIO(output)).slides[0]
    text = " ".join(shape.text for shape in slide.shapes if shape.has_text_frame)
    assert "위협 시나리오" in text and "프롬프트 인젝션" in text
    assert len(slide.shapes) >= 17  # title/header plus a four-column table


def test_html_template_fills_a_dark_callout_with_the_generated_summary():
    template = SimpleNamespace(config=SimpleNamespace(htmlSlides=[
        '<div data-object="true" data-object-type="shape" style="position:absolute;left:0;top:0;width:1920px;height:1080px;background:#FFFFFF"></div>'
        '<div data-object="true" data-object-type="shape" style="position:absolute;left:300px;top:400px;width:1300px;height:400px;background:#1A1A1A"></div>'
    ]))
    output = PPTXGenerator(template).generate(_presentation(_slide(
        "QUOTE", "Decision", {"heading": "Decision", "body": "Ship behind a safety gate.", "bullets": ["One", "Two", "Three"]},
    )))

    slide = Presentation(BytesIO(output)).slides[0]
    filled = next(shape for shape in slide.shapes if shape.has_text_frame and "Ship behind" in shape.text)
    assert _rgb(filled.text_frame.paragraphs[0].runs[0].font.color) == "FFFFFF"


def test_html_template_renders_chart_data_for_chart_slides():
    output = PPTXGenerator(SimpleNamespace(config=SimpleNamespace(htmlSlides=[
        '<div data-object="true" data-object-type="shape" style="position:absolute;left:0;top:0;width:1920px;height:1080px;background:#FFFFFF"></div>'
    ]))).generate(_presentation(_slide(
        "CHART", "ASR", {"heading": "ASR", "chart": {"labels": ["Before", "After"], "values": [48, 11]}},
    )))

    assert any(shape.has_chart for shape in Presentation(BytesIO(output)).slides[0].shapes)


def test_html_template_renders_table_data_for_table_slides():
    output = PPTXGenerator(SimpleNamespace(config=SimpleNamespace(htmlSlides=[
        '<div data-object="true" data-object-type="shape" style="position:absolute;left:0;top:0;width:1920px;height:1080px;background:#FFFFFF"></div>'
    ]))).generate(_presentation(_slide(
        "TABLE", "실적", {"heading": "실적", "table": {"headers": ["기간", "실적"], "rows": [["7/20-7/24", "완료"]]}},
    )))

    slide = Presentation(BytesIO(output)).slides[0]
    texts = [run.text for run in _runs(slide)]
    assert "기간" in texts and "실적" in texts and "완료" in texts


def test_table_with_many_rows_in_small_slot_does_not_overflow():
    # Regression test for a row-height floor overflow bug (formerly at line ~831
    # in _add_table): `row_height = max(h / (len(rows) + 1), 0.3)` applied a 0.3"
    # minimum that, for many rows in a small slot, pushed the table's bottom edge
    # well past the slot's bottom edge. Fixed by removing the floor so row height
    # is always exactly h / (len(rows) + 1).
    #
    # The slot is 400x400px (~2.78"x2.78") so that _add_table's own unrelated,
    # pre-existing `max(slot_h - 0.7, 1.5)` floor (shared with _add_chart, out of
    # scope for this fix) never binds here — this isolates the row-height bug.
    output = PPTXGenerator(SimpleNamespace(config=SimpleNamespace(htmlSlides=[
        '<div data-object="true" data-object-type="shape" style="position:absolute;left:100px;top:100px;width:400px;height:400px;background:#FFFFFF"></div>'
    ]))).generate(_presentation(_slide(
        "TABLE", "Large Table in Small Slot",
        {"heading": "Large Table", "table": {
            "headers": ["Row", "Data"],
            "rows": [[f"Row {i + 1}", f"Value {i + 1}"] for i in range(10)],
        }},
    )))

    slide = Presentation(BytesIO(output)).slides[0]
    rectangles = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE.RECTANGLE]

    # The content slot itself is one rectangle (the largest, by area); every other
    # rectangle is a table cell drawn by _add_table_row (11 rows x 2 columns = 22).
    slot = max(rectangles, key=lambda shape: shape.width * shape.height)
    cells = [shape for shape in rectangles if shape is not slot]
    assert len(cells) == 22, f"expected 22 table-cell rectangles (11 rows x 2 columns), got {len(cells)}"

    slot_bottom = slot.top + slot.height
    max_cell_bottom = max(cell.top + cell.height for cell in cells)
    assert max_cell_bottom <= slot_bottom, (
        f"table overflowed the slot by {(max_cell_bottom - slot_bottom) / 914400:.4f}\" "
        f"(slot bottom {slot_bottom}, deepest cell bottom {max_cell_bottom})"
    )


def test_html_template_without_font_family_uses_default_font():
    template = SimpleNamespace(config=SimpleNamespace(
        htmlTemplate='<div data-object="true" data-object-type="textbox" style="position:absolute;left:120px;top:120px;width:1200px;height:180px;font-size:48px">Title</div>',
    ))
    output = PPTXGenerator(template).generate(_presentation(_slide("TITLE", "Generated", {"heading": "Generated"})))

    assert Presentation(BytesIO(output)).slides[0].shapes[0].text == "Generated"


def test_html_template_without_data_object_markup_falls_back_to_generic_layout():
    # A real-world upload (Genspark export, Tailwind/CSS-class deck) never carries JaSlide's
    # own data-object markers or absolute-pixel inline positioning. Before this fix, such a
    # template produced a fully blank slide instead of the generated content below.
    template = SimpleNamespace(config=SimpleNamespace(
        htmlTemplate=(
            '<style>.card { background: #123456; }</style>'
            '<div class="card"><h1 style="color:#FFFFFF;font-size:56px;font-family:Newsreader">Cover</h1>'
            '<p style="font-size:20px;font-family:Pretendard">Body copy</p></div>'
        ),
        htmlSlides=[
            '<style>.card { background: #123456; }</style>'
            '<div class="card"><h1 style="color:#FFFFFF;font-size:56px;font-family:Newsreader">Cover</h1>'
            '<p style="font-size:20px;font-family:Pretendard">Body copy</p></div>'
        ],
    ))

    pptx = PPTXGenerator(template).generate(
        _presentation(_slide("CONTENT", "Generated title", {"heading": "Generated title", "body": "Generated body"}))
    )

    slide = Presentation(BytesIO(pptx)).slides[0]
    # Template background/fonts are still picked up even without data-object markup.
    assert _rgb(slide.background.fill.fore_color) == "123456"
    texts = [run.text for shape in slide.shapes if shape.has_text_frame
             for paragraph in shape.text_frame.paragraphs for run in paragraph.runs]
    assert "Generated title" in texts
    assert any("Newsreader" == shape.text_frame.paragraphs[0].runs[0].font.name
               for shape in slide.shapes if shape.has_text_frame and shape.text_frame.paragraphs[0].runs)


def test_reusing_a_generator_does_not_accumulate_prior_slides():
    generator = PPTXGenerator()
    presentation = _presentation(_slide("TITLE", "One", {"heading": "One"}))

    generator.generate(presentation)
    second_output = generator.generate(presentation)

    assert len(Presentation(BytesIO(second_output)).slides) == 1


def test_pptx_template_keeps_native_text_and_table_objects_editable():
    source = Presentation()
    slide = source.slides.add_slide(source.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    title.text = "Original title"
    table_shape = slide.shapes.add_table(1, 1, Inches(1), Inches(2), Inches(4), Inches(1))
    table_shape.table.cell(0, 0).text = "Original cell"
    source_buffer = BytesIO()
    source.save(source_buffer)

    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(source_buffer.getvalue()).decode("ascii")))
    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {
        "objectEdits": [
            {"slide": 0, "objectId": str(title.shape_id), "text": "Updated title", "left": 240, "top": 180, "width": 720, "height": 120, "fontSize": 28, "color": "#112233", "bold": True, "italic": True},
            {"slide": 0, "objectId": str(table_shape.shape_id), "cells": [["Updated cell"]]},
        ],
    })))

    generated = Presentation(BytesIO(output))
    assert generated.slides[0].shapes[0].has_text_frame
    assert generated.slides[0].shapes[0].text == "Updated title"
    assert generated.slides[0].shapes[0].left == generated.slide_width * 240 // 1920
    assert generated.slides[0].shapes[0].top == generated.slide_height * 180 // 1080
    run = generated.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
    assert run.font.size == Pt(28) and str(run.font.color.rgb) == "112233" and run.font.bold and run.font.italic
    assert generated.slides[0].shapes[1].has_table
    assert generated.slides[0].shapes[1].table.cell(0, 0).text == "Updated cell"


def test_fit_font_scale_shrinks_only_when_the_text_no_longer_fits():
    # One short line in a roomy box: leave it alone.
    assert fit_font_scale(["짧은 제목"], width_pt=400, height_pt=60, font_pt=24) == 1.0
    # The same box, a title long enough to wrap several times: shrink it.
    long_title = "2026년 3분기 AI 엔지니어링 업무보고 (2026.07.20 ~ 2026.07.24) 추진 실적과 다음 분기 계획"
    assert fit_font_scale([long_title], width_pt=400, height_pt=60, font_pt=24) < 1.0
    # Degenerate boxes must not divide by zero or loop.
    assert fit_font_scale([], 400, 60, 24) == 1.0
    assert fit_font_scale(["글"], 0, 0, 24) == 1.0


def test_pptx_template_shrinks_generated_text_that_outgrew_its_box():
    source = Presentation()
    slide = source.slides.add_slide(source.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(4), Inches(1), Inches(3), Inches(0.5))
    title.text = "Short"
    buffer = BytesIO(); source.save(buffer)
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))
    overflowing = "2026년 3분기 AI 엔지니어링 업무보고 (2026.07.20 ~ 2026.07.24) 추진 실적과 다음 분기 계획"

    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {
        "objectEdits": [{"slide": 0, "objectId": str(title.shape_id), "text": overflowing, "fontSize": 28}],
    })))

    shape = Presentation(BytesIO(output)).slides[0].shapes[0]
    autofit = shape.text_frame._txBody.bodyPr.find(qn("a:normAutofit"))
    # PowerPoint reads fontScale; LibreOffice, which renders the editor preview,
    # only needs the element to be present.
    assert autofit is not None
    assert 0 < int(autofit.get("fontScale")) < 100000
    assert shape.text_frame.word_wrap is True


def test_pptx_template_leaves_text_that_still_fits_at_full_size():
    source = Presentation()
    slide = source.slides.add_slide(source.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    title.text = "Short"
    buffer = BytesIO(); source.save(buffer)
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))

    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {
        "objectEdits": [{"slide": 0, "objectId": str(title.shape_id), "text": "짧은 제목", "fontSize": 18}],
    })))

    autofit = Presentation(BytesIO(output)).slides[0].shapes[0].text_frame._txBody.bodyPr.find(qn("a:normAutofit"))
    assert autofit is not None and autofit.get("fontScale") is None


def test_pptx_template_applies_native_shape_colors():
    source = Presentation()
    slide = source.slides.add_slide(source.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))
    buffer = BytesIO(); source.save(buffer)
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))

    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {"objectEdits": [{"slide": 0, "objectId": str(shape.shape_id), "fillColor": "#112233", "lineColor": "#445566", "lineWidth": 3}]})))

    generated = Presentation(BytesIO(output)).slides[0].shapes[0]
    assert str(generated.fill.fore_color.rgb) == "112233"
    assert str(generated.line.color.rgb) == "445566"
    assert generated.line.width == 3 * 12700


def test_pptx_template_deletes_a_native_object():
    source = Presentation()
    slide = source.slides.add_slide(source.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))
    buffer = BytesIO(); source.save(buffer)
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))

    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {"objectEdits": [{"slide": 0, "objectId": str(shape.shape_id), "delete": True}]})))

    assert len(Presentation(BytesIO(output)).slides[0].shapes) == 0


def test_pptx_template_adds_a_native_image():
    source = Presentation()
    source.slides.add_slide(source.slide_layouts[6])
    buffer = BytesIO(); source.save(buffer)
    image = "data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVQIHWP4z8DwHwAFgAI/ScL1nQAAAABJRU5ErkJggg=="
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))

    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {"objectEdits": [{"slide": 0, "objectId": "new-image", "imageData": image, "left": 200, "top": 100, "width": 400, "height": 200}]})))

    generated = Presentation(BytesIO(output)).slides[0].shapes[0]
    assert generated.shape_type == 13 and generated.left == Presentation(BytesIO(output)).slide_width * 200 // 1920


def test_pptx_template_adds_editable_native_text():
    source = Presentation(); source.slides.add_slide(source.slide_layouts[6])
    buffer = BytesIO(); source.save(buffer)
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))

    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {"objectEdits": [{"slide": 0, "objectId": "new-text", "addText": "Initial", "text": "Edited", "fontSize": 24, "bold": True}]})))

    generated = Presentation(BytesIO(output)).slides[0].shapes[0]
    assert generated.text == "Edited" and generated.text_frame.paragraphs[0].runs[0].font.size == Pt(24) and generated.text_frame.paragraphs[0].runs[0].font.bold


def test_pptx_template_adds_native_shapes_and_lines():
    source = Presentation(); source.slides.add_slide(source.slide_layouts[6])
    buffer = BytesIO(); source.save(buffer)
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))

    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {"objectEdits": [
        {"slide": 0, "objectId": "triangle", "addShape": "triangle", "left": 100, "top": 100, "width": 400, "height": 300, "fillColor": "#FF0000"},
        {"slide": 0, "objectId": "line", "addLine": "straightLine", "left": 600, "top": 200, "width": 500, "height": 80, "lineColor": "#00AA00", "lineWidth": 3},
    ]})))

    shapes = Presentation(BytesIO(output)).slides[0].shapes
    assert len(shapes) == 2 and shapes[0].auto_shape_type == MSO_SHAPE.ISOSCELES_TRIANGLE and shapes[1].shape_type == 9


def test_a_manually_edited_generic_scene_reaches_the_export_instead_of_regenerated_bullets():
    # generation.service.ts's default (no PPTX/HTML source) slides ship as heading/bullets;
    # once GetScene/SaveScene records a manual edit as content.scene, export must place
    # that scene's own objects rather than silently regenerating the heading/bullets layout.
    scene = {
        "width": 1920, "height": 1080,
        "objects": [{
            "id": "manual-text", "type": "text", "x": 77, "y": 88, "width": 400, "height": 100,
            "rotation": 0, "paragraphs": [{"runs": [{"text": "saved edit", "fontSize": 24}]}],
        }],
    }
    output = PPTXGenerator().generate(_presentation(_slide("CONTENT", "Generated heading", {
        "heading": "Generated heading", "bullets": [{"text": "first point", "level": 0}], "scene": scene,
    })))

    slide = Presentation(BytesIO(output)).slides[0]
    assert len(slide.shapes) == 1
    shape = slide.shapes[0]
    assert shape.text_frame.paragraphs[0].runs[0].text == "saved edit"
    assert shape.left == Presentation(BytesIO(output)).slide_width * 77 // 1920


def test_pptx_template_preserves_unedited_table_and_shape():
    source = Presentation(); slide = source.slides.add_slide(source.slide_layouts[6])
    table = slide.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(3), Inches(1)); table.table.cell(0, 0).text = "Keep"
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(1), Inches(2), Inches(1)); shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x11, 0x22, 0x33)
    buffer = BytesIO(); source.save(buffer)
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))
    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {"objectEdits": [{"slide": 0, "objectId": str(table.shape_id), "cells": [["Updated"]]}]})))
    generated = Presentation(BytesIO(output)).slides[0]
    assert generated.shapes[0].table.cell(0, 0).text == "Updated" and str(generated.shapes[1].fill.fore_color.rgb) == "112233"


def test_native_edit_writes_mixed_run_formatting_within_one_paragraph():
    # A user selecting half a sentence and hitting bold must not flatten the
    # rest of it. python-pptx's `shape.text = ...` does exactly that, which is
    # why character-level formatting needs its own path: multiple runs per
    # paragraph, each keeping its own bold/italic/underline/color/size/font.
    source = Presentation(); slide = source.slides.add_slide(source.slide_layouts[6])
    text = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)); text.text = "Old"
    buffer = BytesIO(); source.save(buffer)
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))

    edit = {
        "slide": 0, "objectId": str(text.shape_id),
        "paragraphs": [{
            "text": "Plain bold", "align": "center",
            "runs": [
                {"text": "Plain "},
                {"text": "bold", "bold": True, "color": "#FF0000", "fontSize": 20, "fontFamily": "Newsreader"},
            ],
        }],
    }
    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {"objectEdits": [edit]})))

    paragraph = Presentation(BytesIO(output)).slides[0].shapes[0].text_frame.paragraphs[0]
    assert [run.text for run in paragraph.runs] == ["Plain ", "bold"]
    assert not paragraph.runs[0].font.bold
    assert paragraph.runs[1].font.bold
    assert str(paragraph.runs[1].font.color.rgb) == "FF0000"
    assert paragraph.runs[1].font.size == Pt(20)
    assert paragraph.runs[1].font.name == "Newsreader"
    assert paragraph.alignment == PP_ALIGN.CENTER


def test_native_table_cell_writes_mixed_run_formatting():
    source = Presentation(); slide = source.slides.add_slide(source.slide_layouts[6])
    table = slide.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(4), Inches(1))
    table.table.cell(0, 0).text = "Old"
    buffer = BytesIO(); source.save(buffer)
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))
    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {"objectEdits": [{
        "slide": 0, "objectId": str(table.shape_id), "cells": [[{"paragraphs": [{
            "text": "Plain bold", "runs": [{"text": "Plain "}, {"text": "bold", "bold": True, "color": "#FF0000", "fontSize": 20}],
        }]}]],
    }]})))

    paragraph = Presentation(BytesIO(output)).slides[0].shapes[0].table.cell(0, 0).text_frame.paragraphs[0]
    assert [run.text for run in paragraph.runs] == ["Plain ", "bold"]
    assert not paragraph.runs[0].font.bold and paragraph.runs[1].font.bold
    assert str(paragraph.runs[1].font.color.rgb) == "FF0000" and paragraph.runs[1].font.size == Pt(20)


def test_native_edit_paragraphs_without_runs_still_write_flat_text():
    # Back-compat: the AI-generation path writes {text, level} with no runs.
    source = Presentation(); slide = source.slides.add_slide(source.slide_layouts[6])
    text = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)); text.text = "Old"
    buffer = BytesIO(); source.save(buffer)
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))

    edit = {"slide": 0, "objectId": str(text.shape_id), "paragraphs": [{"text": "New title", "level": 0}]}
    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {"objectEdits": [edit]})))

    assert Presentation(BytesIO(output)).slides[0].shapes[0].text == "New title"


def test_pptx_template_preserves_paragraph_indent_when_text_changes():
    source = Presentation(); slide = source.slides.add_slide(source.slide_layouts[6])
    text = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2)); text.text = "First"; text.text_frame.add_paragraph().text = "Nested"; text.text_frame.paragraphs[1].level = 1
    buffer = BytesIO(); source.save(buffer)
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))
    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {"objectEdits": [{"slide": 0, "objectId": str(text.shape_id), "text": "Changed\nStill nested"}]})))
    assert Presentation(BytesIO(output)).slides[0].shapes[0].text_frame.paragraphs[1].level == 1


def test_pptx_text_replace_preserves_paragraph_alignment():
    source = Presentation(); slide = source.slides.add_slide(source.slide_layouts[6])
    text = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    text.text = "Original"
    text.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    buffer = BytesIO(); source.save(buffer)
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))

    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {
        "objectEdits": [{"slide": 0, "objectId": str(text.shape_id), "text": "Changed"}],
    })))

    paragraph = Presentation(BytesIO(output)).slides[0].shapes[0].text_frame.paragraphs[0]
    assert paragraph.text == "Changed"
    assert paragraph.alignment == PP_ALIGN.CENTER


def test_pptx_table_cell_replace_preserves_paragraph_alignment():
    source = Presentation(); slide = source.slides.add_slide(source.slide_layouts[6])
    table = slide.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(4), Inches(2))
    cell = table.table.cell(0, 0); cell.text = "Original"
    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    buffer = BytesIO(); source.save(buffer)
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))

    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {
        "objectEdits": [{"slide": 0, "objectId": str(table.shape_id), "cells": [["Updated"]]}],
    })))

    paragraph = Presentation(BytesIO(output)).slides[0].shapes[0].table.cell(0, 0).text_frame.paragraphs[0]
    assert paragraph.text == "Updated"
    assert paragraph.alignment == PP_ALIGN.RIGHT


def test_sourcepptx_clones_the_template_slide_when_generated_slides_outnumber_it():
    source = Presentation()
    slide = source.slides.add_slide(source.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    title.text = "Original"
    keep = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(4), Inches(1))
    keep.text = "Keep"
    buffer = BytesIO(); source.save(buffer)
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))

    output = PPTXGenerator(template).generate(_presentation(
        _slide("CONTENT", "", {"objectEdits": [{"slide": 0, "objectId": str(title.shape_id), "text": "First"}]}),
        _slide("CONTENT", "", {"objectEdits": [{"slide": 0, "objectId": str(title.shape_id), "text": "Second"}]}),
    ))

    generated = Presentation(BytesIO(output))
    assert len(generated.slides) == 2
    first_texts = {shape.text for shape in generated.slides[0].shapes if shape.has_text_frame}
    second_texts = {shape.text for shape in generated.slides[1].shapes if shape.has_text_frame}
    assert first_texts == {"First", "Keep"}
    assert second_texts == {"Second", "Keep"}


def test_sourcepptx_preview_for_a_later_slide_returns_only_that_slides_own_edit():
    source = Presentation()
    slide = source.slides.add_slide(source.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    title.text = "Original"
    buffer = BytesIO(); source.save(buffer)
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))

    presentation = _presentation(
        _slide("CONTENT", "", {"objectEdits": [{"slide": 0, "objectId": str(title.shape_id), "text": "First"}]}),
        _slide("CONTENT", "", {"objectEdits": [{"slide": 0, "objectId": str(title.shape_id), "text": "Second"}]}),
    )

    output = PPTXGenerator(template).generate(presentation, slide_index=0)

    generated = Presentation(BytesIO(output))
    assert len(generated.slides) == 1
    assert generated.slides[0].shapes[0].text == "First"


def _two_slide_template():
    source = Presentation()
    for label in ("First template", "Second template"):
        slide = source.slides.add_slide(source.slide_layouts[6])
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = label
    buffer = BytesIO(); source.save(buffer)
    return SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))


def _slide_texts(slide):
    return [shape.text for shape in slide.shapes if shape.has_text_frame]


def test_sourcepptx_renders_the_slides_own_template_index_without_object_edits():
    # A picture-only template slide yields no editable objects, so the slide
    # carries just templateIndex. Falling back to slide 0 rendered (and cloned)
    # the wrong template slide, so the deck lost its real layouts.
    output = PPTXGenerator(_two_slide_template()).generate(_presentation(_slide("CONTENT", "", {"templateIndex": 1})))

    generated = Presentation(BytesIO(output))
    assert len(generated.slides) == 1
    assert _slide_texts(generated.slides[0]) == ["Second template"]


def test_sourcepptx_keeps_only_the_generated_slides_in_their_generated_order():
    # The whole template deck used to be shipped alongside the generated
    # slides, so a 2-slide template plus 2 generated slides exported 4 slides
    # with the originals' untouched content mixed in.
    output = PPTXGenerator(_two_slide_template()).generate(_presentation(
        _slide("CONTENT", "", {"templateIndex": 1}),
        _slide("CONTENT", "", {"templateIndex": 0}),
    ))

    generated = Presentation(BytesIO(output))
    assert [_slide_texts(slide) for slide in generated.slides] == [["Second template"], ["First template"]]


def test_every_preset_shape_name_exports_as_that_shape_not_a_rectangle():
    # The picker sends OOXML preset names. A hand-written lookup covered 20 of
    # them and quietly turned the other 160 into rectangles on export.
    from pptx.enum.shapes import MSO_SHAPE as PRESETS

    source = Presentation(); source.slides.add_slide(source.slide_layouts[6])
    buffer = BytesIO(); source.save(buffer)
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))
    kinds = ["roundRect", "snip2DiagRect", "octagon", "star24", "wedgeEllipseCallout",
             "flowChartMagneticDrum", "curvedUpArrow", "mathNotEqual", "irregularSeal2", "cloudCallout"]

    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {
        "objectEdits": [{"slide": 0, "objectId": f"new-{kind}", "addShape": kind} for kind in kinds],
    })))

    generated = Presentation(BytesIO(output)).slides[0].shapes
    assert [shape.auto_shape_type for shape in generated] == [PRESETS.from_xml(kind) for kind in kinds]


def test_an_inserted_shape_keeps_the_text_typed_into_it_on_the_canvas():
    # Double-clicking a shape in the editor opens the same inline text editor a
    # text box gets, so the shape's edit arrives carrying both addShape and text.
    source = Presentation(); source.slides.add_slide(source.slide_layouts[6])
    buffer = BytesIO(); source.save(buffer)
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))

    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {
        "objectEdits": [{"slide": 0, "objectId": "new-shape", "addShape": "roundRect", "text": "핵심 지표"}],
    })))

    shape = Presentation(BytesIO(output)).slides[0].shapes[0]
    assert shape.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE and shape.text_frame.text == "핵심 지표"


def test_a_table_can_be_inserted_into_a_pptx_slide_and_filled():
    # The toolbar's 표 button only rewrote content.html, so on a PPTX-backed slide —
    # the whole point of a report template — inserting a table was impossible.
    source = Presentation(); source.slides.add_slide(source.slide_layouts[6])
    buffer = BytesIO(); source.save(buffer)
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))

    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {
        "objectEdits": [{
            "slide": 0, "objectId": "new-table", "addTable": {"rows": 2, "columns": 3},
            "left": 200, "top": 300, "width": 900, "height": 300,
            "cells": [["항목", "실적", "계획"], ["NL2SQL", "테스트 완료", "배포"]],
        }],
    })))

    shape = Presentation(BytesIO(output)).slides[0].shapes[0]
    assert shape.has_table
    assert len(shape.table.rows) == 2 and len(shape.table.columns) == 3
    assert [cell.text for cell in shape.table.rows[1].cells] == ["NL2SQL", "테스트 완료", "배포"]


def test_duplicating_an_object_copies_its_look_and_places_the_copy():
    source = Presentation(); slide = source.slides.add_slide(source.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))
    shape.text = "원본"
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x11, 0x22, 0x33)
    buffer = BytesIO(); source.save(buffer)
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))

    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {
        "objectEdits": [{
            "slide": 0, "objectId": "new-copy", "duplicate": str(shape.shape_id),
            "left": 400, "top": 300,
        }],
    })))

    shapes = Presentation(BytesIO(output)).slides[0].shapes
    assert len(shapes) == 2
    copied = shapes[1]
    # Same look and text as the original, but its own id and position.
    assert copied.text_frame.text == "원본"
    assert str(copied.fill.fore_color.rgb) == "112233"
    assert copied.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
    assert copied.shape_id != shape.shape_id
    assert copied.left == Presentation(BytesIO(output)).slide_width * 400 // 1920


def test_an_object_can_be_rotated():
    # 145 insertable shapes include every arrow direction; without rotation an
    # inserted arrow only ever points the way its preset does.
    source = Presentation(); slide = source.slides.add_slide(source.slide_layouts[6])
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(1), Inches(1), Inches(2), Inches(1))
    buffer = BytesIO(); source.save(buffer)
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))

    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {
        "objectEdits": [
            {"slide": 0, "objectId": str(shape.shape_id), "rotation": 135},
            {"slide": 0, "objectId": "new-turned", "addShape": "leftArrow", "rotation": 420},
        ],
    })))

    generated = Presentation(BytesIO(output)).slides[0].shapes
    assert generated[0].rotation == 135
    assert generated[1].rotation == 60, "a rotation beyond one turn wraps"


def test_an_object_can_be_pulled_in_front_of_or_behind_the_others():
    # Overlapping objects are normal on a slide; without this an object could be
    # moved but never pulled out from under the one covering it.
    source = Presentation(); slide = source.slides.add_slide(source.slide_layouts[6])
    first = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))
    middle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2), Inches(1), Inches(2), Inches(1))
    slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(3), Inches(1), Inches(2), Inches(1))
    buffer = BytesIO(); source.save(buffer)
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))

    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {
        "objectEdits": [
            {"slide": 0, "objectId": str(middle.shape_id), "order": "front"},
            {"slide": 0, "objectId": str(first.shape_id), "order": "back"},
        ],
    })))

    shapes = Presentation(BytesIO(output)).slides[0].shapes
    assert shapes[0].shape_id == first.shape_id, "sent to back must be drawn first"
    assert shapes[-1].shape_id == middle.shape_id, "brought to front must be drawn last"
    assert len(shapes) == 3


def test_unknown_shape_names_fall_back_to_a_rectangle_instead_of_failing():
    source = Presentation(); source.slides.add_slide(source.slide_layouts[6])
    buffer = BytesIO(); source.save(buffer)
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))

    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {
        "objectEdits": [{"slide": 0, "objectId": "new-shape", "addShape": "notAPreset"}],
    })))

    assert Presentation(BytesIO(output)).slides[0].shapes[0].auto_shape_type == MSO_SHAPE.RECTANGLE


def test_line_kinds_carry_their_arrowheads_and_dash_pattern():
    source = Presentation(); source.slides.add_slide(source.slide_layouts[6])
    buffer = BytesIO(); source.save(buffer)
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))

    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {
        "objectEdits": [{"slide": 0, "objectId": f"new-{kind}", "addLine": kind}
                        for kind in ("straightLine", "arrowLine", "doubleArrowLine", "dashedLine", "dottedLine")],
    })))

    lines = [shape.line._get_or_add_ln() for shape in Presentation(BytesIO(output)).slides[0].shapes]
    plain, single, double, dashed, dotted = lines
    assert plain.find(qn("a:tailEnd")) is None and plain.find(qn("a:prstDash")) is None
    assert single.find(qn("a:tailEnd")) is not None and single.find(qn("a:headEnd")) is None
    assert double.find(qn("a:tailEnd")) is not None and double.find(qn("a:headEnd")) is not None
    assert dashed.find(qn("a:prstDash")).get("val") == "dash"
    assert dotted.find(qn("a:prstDash")).get("val") == "sysDot"


def test_pptx_table_edits_survive_a_theme_colored_source_cell():
    # Many real-world decks color table text via a theme/scheme reference
    # (design-system driven) instead of an explicit RGB value. Reading
    # `.rgb` off such a run raises AttributeError, which must not crash the
    # whole render — the edit should still land, just without copying color.
    source = Presentation(); slide = source.slides.add_slide(source.slide_layouts[6])
    table = slide.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(4), Inches(2))
    cell = table.table.cell(0, 0); cell.text = "Original"
    cell.text_frame.paragraphs[0].runs[0].font.color.theme_color = MSO_THEME_COLOR.ACCENT_1
    buffer = BytesIO(); source.save(buffer)
    template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))

    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {
        "objectEdits": [{"slide": 0, "objectId": str(table.shape_id), "cells": [["Updated"]]}],
    })))

    assert Presentation(BytesIO(output)).slides[0].shapes[0].table.cell(0, 0).text == "Updated"


def test_pptx_table_edits_keep_cell_text_style():
    source = Presentation(); slide = source.slides.add_slide(source.slide_layouts[6]); table = slide.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(4), Inches(2))
    cell = table.table.cell(0, 0); cell.text = "Original"; run = cell.text_frame.paragraphs[0].runs[0]; run.font.size = Pt(18); run.font.bold = True
    buffer = BytesIO(); source.save(buffer); template = SimpleNamespace(config=SimpleNamespace(sourcePptx=base64.b64encode(buffer.getvalue()).decode("ascii")))
    output = PPTXGenerator(template).generate(_presentation(_slide("CONTENT", "", {"objectEdits": [{"slide": 0, "objectId": str(table.shape_id), "cells": [["Updated"]]}]})))
    updated = Presentation(BytesIO(output)).slides[0].shapes[0].table.cell(0, 0).text_frame.paragraphs[0].runs[0]
    assert updated.font.size == Pt(18) and updated.font.bold
