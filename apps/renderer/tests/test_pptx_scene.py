from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from apps.renderer.src.services.pptx_scene import apply_edits_to_scene, pptx_to_scene


def _deck_bytes(build) -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    build(slide)
    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_a_text_box_keeps_its_shape_id_and_per_run_formatting():
    def build(slide):
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        paragraph = box.text_frame.paragraphs[0]
        bold_run = paragraph.add_run()
        bold_run.text = "Plain "
        plain_run = paragraph.add_run()
        plain_run.text = "bold"
        plain_run.font.bold = True
        plain_run.font.size = Pt(20)
        plain_run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)

    scene = pptx_to_scene(_deck_bytes(build))["slides"][0]
    text = next(item for item in scene["objects"] if item["type"] == "text")

    assert text["sourceRef"] == {"kind": "pptx-shape", "shapeId": text["id"]}
    runs = text["paragraphs"][0]["runs"]
    assert [run["text"] for run in runs] == ["Plain ", "bold"]
    assert not runs[0].get("bold")
    assert runs[1]["bold"] is True
    assert runs[1]["fontSize"] == 20
    assert runs[1]["color"] == "#FF0000"


def test_a_table_cell_keeps_its_fill_and_paragraphs():
    def build(slide):
        table = slide.shapes.add_table(1, 2, Inches(1), Inches(1), Inches(6), Inches(1)).table
        table.cell(0, 0).text = "Header"
        table.cell(0, 0).fill.solid()
        table.cell(0, 0).fill.fore_color.rgb = RGBColor(0xD9, 0xD9, 0xD9)

    scene = pptx_to_scene(_deck_bytes(build))["slides"][0]
    table = next(item for item in scene["objects"] if item["type"] == "table")

    assert table["cells"][0][0]["paragraphs"][0]["runs"][0]["text"] == "Header"
    assert table["cells"][0][0]["fill"] == "#D9D9D9"
    assert len(table["rowHeights"]) == 1 and len(table["columnWidths"]) == 2


def test_a_table_keeps_cell_borders_and_merge_geometry():
    def add_border(cell, side, color, width):
        tc_pr = cell._tc.get_or_add_tcPr()
        line = OxmlElement(f"a:ln{side}")
        line.set("w", str(width))
        fill = OxmlElement("a:solidFill")
        rgb = OxmlElement("a:srgbClr")
        rgb.set("val", color)
        fill.append(rgb)
        line.append(fill)
        tc_pr.append(line)

    def build(slide):
        table = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(6), Inches(2)).table
        table.cell(0, 0).text = "Merged header"
        table.cell(0, 0).merge(table.cell(0, 1))
        add_border(table.cell(0, 0), "B", "FF0000", 25400)

    scene = pptx_to_scene(_deck_bytes(build))["slides"][0]
    table = next(item for item in scene["objects"] if item["type"] == "table")
    header = table["cells"][0][0]

    assert header["colSpan"] == 2
    assert header["border"]["bottom"] == {"color": "#FF0000", "width": 2}
    assert table["cells"][0][1]["spanned"] is True


def test_a_table_keeps_a_theme_coloured_border():
    def build(slide):
        cell = slide.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(4), Inches(1)).table.cell(0, 0)
        line = OxmlElement("a:lnT")
        fill = OxmlElement("a:solidFill")
        color = OxmlElement("a:schemeClr")
        color.set("val", "accent1")
        fill.append(color)
        line.append(fill)
        cell._tc.get_or_add_tcPr().append(line)

    scene = pptx_to_scene(_deck_bytes(build))["slides"][0]
    cell = next(item for item in scene["objects"] if item["type"] == "table")["cells"][0][0]

    assert cell["border"]["top"] == {"color": "#4F81BD", "width": 1}


def test_a_table_keeps_cell_padding_and_vertical_alignment():
    def build(slide):
        cell = slide.shapes.add_table(1, 1, Inches(1), Inches(1), Inches(4), Inches(2)).table.cell(0, 0)
        cell.text = "Centered"
        cell.margin_top = Pt(10)
        cell.margin_right = Pt(12)
        cell.margin_bottom = Pt(8)
        cell.margin_left = Pt(6)
        cell.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    scene = pptx_to_scene(_deck_bytes(build))["slides"][0]
    cell = next(item for item in scene["objects"] if item["type"] == "table")["cells"][0][0]

    assert cell["verticalAlign"] == "middle"
    assert cell["padding"] == {"top": 20, "right": 24, "bottom": 16, "left": 12}


def test_an_autoshape_keeps_its_preset_geometry_fill_and_rotation():
    def build(slide):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x25, 0x63, 0xEB)
        shape.rotation = 15

    scene = pptx_to_scene(_deck_bytes(build))["slides"][0]
    shape_object = next(item for item in scene["objects"] if item["type"] == "shape")

    # The exact OOXML preset name, not a lossy enum round-trip: _preset_shape on
    # export feeds this straight to MSO_SHAPE.from_xml, so keeping the deck's own
    # string is what makes re-exporting produce the same silhouette it started with.
    assert shape_object["shape"] == "roundRect"
    assert shape_object["fill"] == "#2563EB"
    assert shape_object["rotation"] == 15


def test_a_straight_connector_is_a_line_not_a_shape():
    def build(slide):
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(1), Inches(4), Inches(1))
        connector.line.color.rgb = RGBColor(0x00, 0xAA, 0x00)
        connector.line.width = Pt(3)

    scene = pptx_to_scene(_deck_bytes(build))["slides"][0]

    assert [item["type"] for item in scene["objects"]] == ["line"]
    line = scene["objects"][0]
    assert line["stroke"] == "#00AA00"
    assert line["lineStyle"] == "straightLine"


def test_rotation_defaults_to_zero_when_the_deck_states_none():
    def build(slide):
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1)).text_frame.paragraphs[0].add_run().text = "x"

    scene = pptx_to_scene(_deck_bytes(build))["slides"][0]

    assert scene["objects"][0]["rotation"] == 0


def test_scene_canvas_size_matches_the_editor_stage():
    scene = pptx_to_scene(_deck_bytes(lambda slide: None))["slides"][0]

    assert scene["width"] == 1920
    assert scene["height"] == 1080


def test_apply_edits_to_scene_moves_an_existing_object():
    def build(slide):
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1)).text_frame.paragraphs[0].add_run().text = "x"

    scene = pptx_to_scene(_deck_bytes(build))["slides"][0]
    shape_id = scene["objects"][0]["id"]

    edited = apply_edits_to_scene(scene, [{"objectId": shape_id, "left": 300, "top": 300, "width": 200, "height": 50}])

    assert (edited["objects"][0]["x"], edited["objects"][0]["y"]) == (300, 300)
    assert (edited["objects"][0]["width"], edited["objects"][0]["height"]) == (200, 50)


def test_apply_edits_to_scene_inserts_a_new_shape_with_no_source_shape():
    scene = pptx_to_scene(_deck_bytes(lambda slide: None))["slides"][0]

    edited = apply_edits_to_scene(scene, [{
        "objectId": "new-shape-1", "left": 10, "top": 20, "width": 100, "height": 50,
        "addShape": "roundRect", "fillColor": "#FF0000", "lineColor": "#000000", "lineWidth": 2,
    }])

    assert len(edited["objects"]) == 1
    inserted = edited["objects"][0]
    assert inserted == {
        "id": "new-shape-1", "x": 10, "y": 20, "width": 100, "height": 50, "rotation": 0,
        "type": "shape", "shape": "roundRect", "fill": "#FF0000", "stroke": "#000000", "strokeWidth": 2,
    }


def test_apply_edits_to_scene_deletes_an_object():
    def build(slide):
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1)).text_frame.paragraphs[0].add_run().text = "x"

    scene = pptx_to_scene(_deck_bytes(build))["slides"][0]
    shape_id = scene["objects"][0]["id"]

    edited = apply_edits_to_scene(scene, [{"objectId": shape_id, "delete": True}])

    assert edited["objects"] == []


def test_apply_edits_to_scene_converts_legacy_string_cells_to_paragraphs():
    # generation.service.ts writes table objectEdits the same way _apply_native_edit
    # (pptx_generator.py) reads them — plain strings per cell, never the scene's own
    # {paragraphs: [...]} shape. The scene view must normalize these, or SceneCanvas's
    # table renderer crashes reading `.paragraphs` off a bare string.
    def build(slide):
        slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(6), Inches(2))

    scene = pptx_to_scene(_deck_bytes(build))["slides"][0]
    table_id = scene["objects"][0]["id"]

    edited = apply_edits_to_scene(scene, [{
        "objectId": table_id,
        "cells": [["추진실적", "추진계획"], ["a", "b"]],
    }])

    table = edited["objects"][0]
    assert table["cells"][0][0]["paragraphs"][0]["runs"][0]["text"] == "추진실적"
    assert table["cells"][1][1]["paragraphs"][0]["runs"][0]["text"] == "b"


def test_apply_edits_to_scene_applies_a_legacy_plain_text_edit():
    # generation.service.ts also writes plain-text objectEdits the same way
    # _apply_native_edit (pptx_generator.py) reads them for a text shape —
    # {"text": "..."} , never the scene's own {paragraphs: [...]} shape. Before
    # this fix, the scene view silently ignored these edits and kept showing
    # the template's stale placeholder text instead of the generated one.
    def build(slide):
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        box.text_frame.paragraphs[0].add_run().text = "stale placeholder"

    scene = pptx_to_scene(_deck_bytes(build))["slides"][0]
    text_id = scene["objects"][0]["id"]

    edited = apply_edits_to_scene(scene, [{"objectId": text_id, "text": "0730 업무보고"}])

    assert edited["objects"][0]["paragraphs"][0]["runs"][0]["text"] == "0730 업무보고"


def test_apply_edits_to_scene_duplicates_an_existing_object():
    def build(slide):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(1), Inches(2), Inches(1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x00, 0xFF, 0x00)

    scene = pptx_to_scene(_deck_bytes(build))["slides"][0]
    original_id = scene["objects"][0]["id"]

    edited = apply_edits_to_scene(scene, [{"objectId": "copy-1", "duplicate": original_id, "left": 500, "top": 500}])

    assert len(edited["objects"]) == 2
    copy = next(item for item in edited["objects"] if item["id"] == "copy-1")
    assert copy["fill"] == "#00FF00"
    assert (copy["x"], copy["y"]) == (500, 500)
    assert "sourceRef" not in copy
